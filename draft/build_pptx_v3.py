"""
build_pptx_v3.py  –  Professional 14-slide presentation with embedded charts
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x0D, 0x2B, 0x55)
TEAL  = RGBColor(0x00, 0x7A, 0x8A)
GOLD  = RGBColor(0xE8, 0xA0, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF5, 0xF6, 0xFA)
MGRAY = RGBColor(0xCC, 0xCC, 0xCC)
RED   = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)

IMG   = '/Users/leekim/prj/deloitte/'
W, H  = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def tb(slide, text, l, t, w, h, size=12, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
    add_rect(slide, 0, 1.1, 13.33, 0.06, GOLD)
    tb(slide, title, 0.4, 0.12, 12, 0.65, size=22, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.4, 0.72, 12, 0.35, size=11, color=GOLD)

def page_bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, LGRAY)

def kpi_box(slide, l, t, w, h, label, value, sub='', val_color=TEAL):
    add_rect(slide, l, t, w, h, WHITE)
    add_rect(slide, l, t, w, 0.06, TEAL)
    tb(slide, value, l+0.1, t+0.15, w-0.2, 0.55, size=26, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    tb(slide, label, l+0.1, t+0.72, w-0.2, 0.32, size=9,  bold=True, color=NAVY,     align=PP_ALIGN.CENTER)
    if sub:
        tb(slide, sub, l+0.1, t+1.02, w-0.2, 0.25, size=8, color=MGRAY, align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, title, bullets, title_bg=TEAL):
    add_rect(slide, l, t, w, h, WHITE)
    add_rect(slide, l, t, w, 0.38, title_bg)
    tb(slide, title, l+0.12, t+0.06, w-0.24, 0.3, size=10, bold=True, color=WHITE)
    y = t + 0.45
    for b in bullets:
        tb(slide, f'• {b}', l+0.12, y, w-0.24, 0.26, size=8.5, color=NAVY)
        y += 0.27

def img(slide, path, l, t, w, h=None):
    full = IMG + path if not path.startswith('/') else path
    if not os.path.exists(full):
        return
    if h:
        slide.shapes.add_picture(full, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(full, Inches(l), Inches(t), width=Inches(w))

def divider(slide, t):
    add_rect(slide, 0.4, t, 12.53, 0.03, TEAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl, 0, 0, 0.35, 7.5, TEAL)
add_rect(sl, 0.35, 3.1, 12.98, 0.07, GOLD)

tb(sl, 'SME Credit Scorecard Development', 1.0, 1.4, 11.5, 1.1, size=36, bold=True, color=WHITE)
tb(sl, 'Credit Risk Model: Weight Assessment, EDA & PD Estimation', 1.0, 2.6, 11.5, 0.55, size=16, color=GOLD)
tb(sl, 'Candidate Technical Presentation  |  Deloitte Credit Risk', 1.0, 3.4, 11.5, 0.45, size=13, color=MGRAY)
tb(sl, 'Dataset: 557 SME Borrowers  |  8 Features  |  50 Defaults (9.0%)', 1.0, 4.1, 11.5, 0.4, size=11, color=MGRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Executive Summary
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'Executive Summary', 'Key findings across weight assessment, model performance and recommendations')

# KPIs row
kpi_box(sl, 0.4,  1.3, 2.3, 1.4, 'AUC-ROC',        '0.90',    'Excellent discrimination')
kpi_box(sl, 2.8,  1.3, 2.3, 1.4, 'Gini Coefficient','0.80',    'Strong predictive power')
kpi_box(sl, 5.2,  1.3, 2.3, 1.4, 'Portfolio Default','9.0%',   '50 / 557 borrowers')
kpi_box(sl, 7.6,  1.3, 2.3, 1.4, 'Top Predictor',   'C3',      'Worst CB Delinquency IV=5.43')
kpi_box(sl, 10.0, 1.3, 2.93, 1.4,'Misweighted Feature','D/E Ratio','LR 22% vs Expert 10%')

divider(sl, 2.95)

# Three-column findings
cols = [
    ('Q1: Weight Assessment',
     ['C3 (CB Delinquency) IV=5.43 dominates',
      'D/E Ratio underweighted: 10% → 22%',
      'DSCR underweighted: 10% → 16%',
      'F3 Sale Growth IV=0.04 — near zero signal']),
    ('Q2: Scorecard Enhancements',
     ['Apply Laplace smoothing to sparse bins',
      'Merge zero-default bins to avoid WoE explosion',
      'Recalibrate PD — class_weight inflates PD 4×',
      'Add macroeconomic overlay (GDP, rates)']),
    ('Q4/Q5: Model Results',
     ['LR with WoE encoding — no scaling required',
      'HL test fails: extreme WoE distorts probabilities',
      'Binomial: predicted 35% vs actual 9% (calibrate)',
      'Score range 300–850 after min-max scaling']),
]
x = 0.4
for title, bullets in cols:
    card(sl, x, 3.1, 4.1, 3.95, title, bullets)
    x += 4.22

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Background & Data Overview
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'Background & Data Overview', 'SME portfolio scorecard — expert-weighted system under review')

add_rect(sl, 0.4, 1.25, 5.8, 5.8, WHITE)
tb(sl, 'Context', 0.6, 1.35, 5.4, 0.35, size=12, bold=True, color=NAVY)
divider(sl, 1.72)
context = [
    ('Dataset',     '557 SME borrowers, 8 features, 9.0% default rate'),
    ('Target',       'Default Flag (1 = default, 0 = performing)'),
    ('Scorecard',    'Expert-weighted, C + F components, score per bin'),
    ('C Component',  '50% weight: C1 Restructuring, C2 Rescheduling, C3 CB Delinquency'),
    ('F Component',  '50% weight: F1–F5 financial ratios equally at 10% each'),
    ('Missing Data', 'C2: 282 NaN (51%), Net Fixed Asset: 26 (5%)'),
    ('Sentinels',    '99999998/99999999 used for negative equity & zero denominators'),
    ('Objective',    'Assess expert weights using IV; develop LR-based PD model'),
]
y = 1.85
for k, v in context:
    tb(sl, k, 0.6, y, 1.7, 0.32, size=9, bold=True, color=TEAL)
    tb(sl, v, 2.3, y, 3.8, 0.32, size=9, color=NAVY)
    y += 0.39

# Feature table right
add_rect(sl, 6.5, 1.25, 6.45, 5.8, WHITE)
tb(sl, 'Feature Summary', 6.7, 1.35, 6.0, 0.35, size=12, bold=True, color=NAVY)
divider(sl, 1.72)
hdrs = ['Code','Feature','Type','IV','Expert Wt']
xs   = [6.55, 7.0, 9.4, 10.8, 11.7]
ws   = [0.4,  2.35, 1.35, 0.85, 0.85]
y = 1.85
for i, h in enumerate(hdrs):
    tb(sl, h, xs[i], y, ws[i], 0.3, size=8.5, bold=True, color=NAVY)
divider(sl, 2.17)
rows = [
    ('C1','Debt Restructuring History','Cat','1.17','12.5%'),
    ('C2','Debt Rescheduling History','Cat','0.57','12.5%'),
    ('C3','Worst CB Delinquency 2yr','Cat','5.43','25.0%'),
    ('F1','Net Fixed Asset','Num','2.60','10.0%'),
    ('F2','Net Profit Margin (%)','Num','1.28','10.0%'),
    ('F3','Sale Growth (%)','Num','0.04','10.0%'),
    ('F4','Debt-to-Equity Ratio','Num','1.46','10.0%'),
    ('F5','Debt Service Coverage Ratio','Num','1.58','10.0%'),
]
y = 2.25
for r in rows:
    bg = LGRAY if rows.index(r) % 2 == 0 else WHITE
    add_rect(sl, 6.52, y-0.03, 6.38, 0.33, bg)
    for i, v in enumerate(r):
        col = RED if v == '0.04' else NAVY
        tb(sl, v, xs[i], y, ws[i], 0.3, size=8, color=col)
    y += 0.35

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 – EDA: Univariate Analysis
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'EDA: Univariate Analysis', 'Distribution and default rates across categorical and numerical features')

img(sl, 'eda_univariate_boxplot.png', 0.4, 1.25, 8.2, 4.5)

add_rect(sl, 8.8, 1.25, 4.15, 4.5, WHITE)
tb(sl, 'Key Observations', 9.0, 1.38, 3.8, 0.35, size=11, bold=True, color=NAVY)
obs = [
    ('C1 – Restructuring', 'Binaries with 6-12 month rescheduling show 44% default rate vs 8% baseline'),
    ('C3 – CB Delinquency', '>90-day overdue (>1yr ago) shows 100% default — strongest single predictor'),
    ('Net Profit Margin', 'Heavy left-skew (-113 to +100%), median near -8%; most SMEs unprofitable'),
    ('Sale Growth', 'Extreme right skew (max +10.95 skewness); many zero-growth observations'),
    ('DSCR', 'Defaulters concentrated at DSCR ≤ 0; non-defaulters spread across positive values'),
]
y = 1.85
for title, detail in obs:
    tb(sl, f'▸ {title}', 9.0, y, 3.7, 0.22, size=8.5, bold=True, color=TEAL)
    tb(sl, detail, 9.0, y+0.22, 3.7, 0.38, size=8, color=NAVY)
    y += 0.67

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 – EDA: Correlation & Multicollinearity
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'EDA: Correlation & Multicollinearity', "Pearson correlation for numerical features; Cramér's V for categorical pairs")

img(sl, 'eda_correlation.png', 0.4, 1.25, 6.5, 4.5)
img(sl, 'eda_cramers_v.png',   7.1, 1.25, 3.6, 3.0)

add_rect(sl, 7.1, 4.4, 5.85, 2.65, WHITE)
tb(sl, 'Multicollinearity Findings', 7.3, 4.52, 5.5, 0.32, size=11, bold=True, color=NAVY)
findings = [
    'Cramér\'s V: C1↔C2=0.40, C1↔C3=0.37, C2↔C3=0.39 — moderate correlation among credit history features',
    'Numerical KS test: Net Fixed Asset (p=0.0005) and DSCR (p=0.0004) are statistically significant separators',
    'Net Profit Margin and Sale Growth show weak correlation with default (p>0.05)',
    'No severe multicollinearity — all features retained; VIF check recommended for production',
]
y = 4.9
for f in findings:
    tb(sl, f'• {f}', 7.3, y, 5.5, 0.42, size=8.5, color=NAVY)
    y += 0.45

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 – EDA: Monotonicity & WoE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'EDA: Monotonicity & WoE Analysis', 'Weight of Evidence validates risk ordering across bins')

img(sl, 'eda_monotonicity.png', 0.4, 1.25, 12.5, 4.8)

add_rect(sl, 0.4, 6.15, 12.5, 1.1, WHITE)
tb(sl, 'Monotonicity Findings:', 0.6, 6.27, 2.5, 0.35, size=10, bold=True, color=NAVY)
pts = [
    ('C1 Restructuring', 'Non-monotone: "Never restructured" has lower WoE than "No history" — bins should be merged'),
    ('C3 CB Delinquency', 'Broadly monotone: longer overdue → higher WoE; zero-default bins produce extreme WoE (±20)'),
    ('F4 D/E Ratio', 'Partially monotone; (0.5,1) bin has zero defaults → WoE = -20.23 (Laplace smoothing needed)'),
]
x = 0.4
for label, detail in pts:
    tb(sl, f'▸ {label}:', x, 6.62, 3.0, 0.22, size=8.5, bold=True, color=TEAL)
    tb(sl, detail, x, 6.85, 4.0, 0.32, size=8, color=NAVY)
    x += 4.2

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Q1: Weight Assessment (IV)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'Q1: Feature Weight Assessment', 'Information Value (IV) as objective measure vs expert-assigned weights')

img(sl, 'iv_comparison_chart.png', 0.4, 1.25, 7.8, 4.0)

add_rect(sl, 8.4, 1.25, 4.55, 5.8, WHITE)
tb(sl, 'Assessment Findings', 8.6, 1.38, 4.2, 0.35, size=11, bold=True, color=NAVY)
divider(sl, 1.76)

rows_iv = [
    ('C3 CB Delinquency', '25%', '5.43', '✅ Aligned'),
    ('C1 Restructuring',  '12.5%', '1.17', '✅ Reasonable'),
    ('C2 Rescheduling',   '12.5%', '0.57', '⚠ Overweighted'),
    ('F1 Net Fixed Asset','10%', '2.60', '⚠ Underweighted'),
    ('F4 D/E Ratio',      '10%', '1.46', '⚠ Underweighted'),
    ('F5 DSCR',           '10%', '1.58', '⚠ Underweighted'),
    ('F2 Profit Margin',  '10%', '1.28', '⚠ Underweighted'),
    ('F3 Sale Growth',    '10%', '0.04', '❌ Weak — remove'),
]
hdrs2 = ['Feature','Expert','IV','Status']
xs2   = [8.45, 10.25, 10.85, 11.45]
ws2   = [1.75, 0.55, 0.55, 1.5]
y = 1.85
for i, h in enumerate(hdrs2):
    tb(sl, h, xs2[i], y, ws2[i], 0.28, size=8.5, bold=True, color=NAVY)
divider(sl, 2.16)
y = 2.22
for r in rows_iv:
    bg = LGRAY if rows_iv.index(r) % 2 == 0 else WHITE
    add_rect(sl, 8.42, y-0.02, 4.48, 0.32, bg)
    for i, v in enumerate(r):
        col = RED if '❌' in v else (GOLD if '⚠' in v else GREEN if '✅' in v else NAVY)
        tb(sl, v, xs2[i], y, ws2[i], 0.28, size=8, color=col)
    y += 0.33

add_rect(sl, 0.4, 5.4, 7.8, 1.65, WHITE)
tb(sl, 'Recommendation: Proposed Weight Rebalancing', 0.6, 5.52, 7.4, 0.3, size=10, bold=True, color=NAVY)
recs = 'Increase C3 → 30% | F1+F4+F5 → 15% each | Remove F3 (IV<0.1) | Rebalance C1/C2 toward 10%/15%'
tb(sl, recs, 0.6, 5.85, 7.4, 0.55, size=9, color=NAVY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Q2: Scorecard Enhancements
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'Q2: Scorecard Enhancement Suggestions', 'Six areas to improve model robustness, calibration and coverage')

enhancements = [
    ('1. Laplace Smoothing on WoE',
     ['Bins with 0 defaults produce WoE = ±20',
      'Apply ε=0.5 smoothing: pct_bad=(bad+ε)/(total_bad+ε×n)',
      'Prevents extreme score swings for sparse bins']),
    ('2. Bin Merging Strategy',
     ['"Restructured in past" has 0 defaults — merge with adjacent',
      'Improves monotonicity and WoE stability',
      'Rule: merge if n < 5 or default rate = 0 with no risk logic']),
    ('3. Remove F3 Sale Growth',
     ['IV = 0.04 → weak predictor (threshold: IV < 0.10)',
      'Chi-square p = 0.89 — not statistically significant',
      'Inclusion adds noise without predictive gain']),
    ('4. Recalibrate PD Probabilities',
     ['class_weight="balanced" inflates predicted PD 4×',
      'Use Platt scaling (CalibratedClassifierCV)',
      'Or set base_odds to actual portfolio default rate']),
    ('5. Macroeconomic Overlay',
     ['Incorporate GDP growth, interest rate, oil price',
      'Lag variables 1–4 quarters (Granger causality test)',
      'Two-stage model: micro scorecard + macro adjustment']),
    ('6. Rating Class Calibration',
     ['Define score bands using log-linear PD spacing',
      'Each class needs ≥20 defaults for stable PD estimate',
      'Validate monotonicity: PD must strictly increase band-by-band']),
]
xs3 = [0.4, 4.55, 8.7]
ys3 = [1.25, 4.05]
idx = 0
for r in range(2):
    for c in range(3):
        title, bullets = enhancements[idx]
        card(sl, xs3[c], ys3[r], 3.9, 2.55, title, bullets)
        idx += 1

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Q3: EDA Methodology
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'Q3: EDA Methodology for Credit Risk Datasets', 'Structured five-step framework for borrower, financial and macroeconomic data')

steps = [
    ('Step 1\nData Quality', ['Missing value audit & sentinel detection (99999998/99999999)',
                               'Skewness check → IQR-based outlier treatment',
                               'Duplicate and data type validation']),
    ('Step 2\nUnivariate', ['Categorical: frequency, default rate per bin, concentration check',
                             'Numerical: histogram, boxplot, describe()',
                             'Flag features with >50% single-bin concentration']),
    ('Step 3\nBivariate', ['Numerical vs Target: KS test, Pearson correlation',
                            'Categorical vs Target: Chi-square, Cramér\'s V',
                            'IV/WoE: quantify predictive power per bin']),
    ('Step 4\nMultivariate', ['Correlation matrix (Pearson + Cramér\'s V)',
                               'VIF for numerical multicollinearity',
                               'PCA if >15 features for dimensionality view']),
    ('Step 5\nMacroeconomic', ['GDP growth, interest rates, oil price vs aggregate default rate',
                                'Granger causality test with 1–4 quarter lags',
                                'Two-stage model: micro + macro components']),
]
x = 0.4
for i, (title, bullets) in enumerate(steps):
    add_rect(sl, x, 1.25, 2.42, 5.5, WHITE)
    add_rect(sl, x, 1.25, 2.42, 0.7, NAVY if i == 0 else TEAL)
    tb(sl, title, x+0.1, 1.3, 2.2, 0.6, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y = 2.05
    for b in bullets:
        tb(sl, f'• {b}', x+0.12, y, 2.2, 0.5, size=8.5, color=NAVY)
        y += 0.55
    # step connector
    if i < 4:
        add_rect(sl, x+2.42, 2.45, 0.15, 0.25, GOLD)
    x += 2.57

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Q4: PD Model Development
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'Q4: PD Model Development — Logistic Regression', 'Theory, assumptions and implementation steps')

# Left: theory
add_rect(sl, 0.4, 1.25, 6.2, 5.8, WHITE)
tb(sl, 'Model Theory', 0.6, 1.38, 5.8, 0.35, size=11, bold=True, color=NAVY)
divider(sl, 1.76)
theory = [
    ('Log-odds formulation',
     'log(p/(1-p)) = β₀ + β₁WoE₁ + … + βₙWoEₙ\nWoE encoding linearises the log-odds relationship'),
    ('MLE estimation',
     'Parameters estimated by maximising log-likelihood\nNo closed-form solution — iterative (Newton-Raphson)'),
    ('WoE encoding effect',
     'Replaces categorical bins with ln(%Bad/%Good)\nMakes LR assumptions valid; handles non-linearity'),
    ('Class imbalance',
     'class_weight="balanced" adjusts loss function\nNote: distorts probability scale — recalibration required'),
    ('Key assumptions',
     'Linear log-odds; no perfect multicollinearity\nObservations independent; large-sample MLE'),
]
y = 1.9
for title, detail in theory:
    tb(sl, f'▸ {title}', 0.6, y, 5.8, 0.24, size=9, bold=True, color=TEAL)
    tb(sl, detail, 0.6, y+0.24, 5.8, 0.44, size=8.5, color=NAVY)
    y += 0.78

# Right: implementation steps
add_rect(sl, 6.8, 1.25, 6.15, 5.8, WHITE)
tb(sl, 'Implementation Steps', 7.0, 1.38, 5.8, 0.35, size=11, bold=True, color=NAVY)
divider(sl, 1.76)
steps6 = [
    ('1. Data Preparation', 'Bin features per scorecard; fill NaN as "Not available"; replace sentinels'),
    ('2. WoE Encoding',     'Compute WoE per bin; apply Laplace smoothing for sparse bins'),
    ('3. Train/Test Split', '80/20 stratified split; EPV check (need ≥10 defaults per feature)'),
    ('4. Model Training',   'sklearn LogisticRegression; class_weight="balanced" for imbalance'),
    ('5. PDO Calibration',  'Factor=PDO/ln2=28.85; Offset=Base_Score−Factor×ln(Base_Odds)'),
    ('6. Score Mapping',    'Score_i = Offset + Factor×β₀ + Σ[−Factor×βᵢ×WoEᵢ]'),
]
y = 1.9
for i, (title, detail) in enumerate(steps6):
    add_rect(sl, 6.85, y-0.03, 0.35, 0.35, TEAL if i % 2 == 0 else NAVY)
    tb(sl, str(i+1), 6.87, y-0.01, 0.3, 0.3, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, title, 7.28, y, 5.5, 0.24, size=9, bold=True, color=TEAL)
    tb(sl, detail, 7.28, y+0.24, 5.5, 0.4, size=8.5, color=NAVY)
    y += 0.78

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Model Results
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'Model Results', 'LR coefficients, score distribution and performance metrics')

# Top KPIs
kpi_box(sl, 0.4,  1.25, 2.3, 1.3, 'AUC-ROC',   '0.90', 'Excellent')
kpi_box(sl, 2.8,  1.25, 2.3, 1.3, 'Gini',       '0.80', 'Strong')
kpi_box(sl, 5.2,  1.25, 2.3, 1.3, 'Accuracy',   '65%',  'Threshold-dependent')
kpi_box(sl, 7.6,  1.25, 2.3, 1.3, 'Recall (D)', '80%',  'Default detection')
kpi_box(sl, 10.0, 1.25, 2.93,1.3, 'Factor/Offset','28.85 / 487.1', 'PDO=20, Base=600')

# LR coef table left
add_rect(sl, 0.4, 2.7, 5.5, 4.55, WHITE)
tb(sl, 'LR Coefficients & Normalised Importance', 0.6, 2.82, 5.1, 0.32, size=10, bold=True, color=NAVY)
coef_rows = [
    ('C3 Worst CB Delinquency','0.897','26.60%','25%'),
    ('F4 Debt-to-Equity Ratio','0.745','22.17%','10%'),
    ('F5 DSCR',               '0.545','16.22%','10%'),
    ('C2 Debt Rescheduling',  '0.469','13.96%','12.5%'),
    ('F1 Net Fixed Asset',    '0.282',' 8.43%','10%'),
    ('C1 Debt Restructuring', '0.238',' 7.09%','12.5%'),
    ('F2 Net Profit Margin',  '0.182',' 5.52%','10%'),
]
xc = [0.45, 2.35, 3.5, 4.5]
wc = [1.85, 1.1, 0.95, 0.9]
hd = ['Feature','Coef','LR Wt%','Expert']
y = 3.17
for i, h in enumerate(hd):
    tb(sl, h, xc[i], y, wc[i], 0.27, size=8.5, bold=True, color=NAVY)
divider(sl, 3.47)
y = 3.52
for r in coef_rows:
    bg = LGRAY if coef_rows.index(r) % 2 == 0 else WHITE
    add_rect(sl, 0.42, y-0.02, 5.44, 0.3, bg)
    for i, v in enumerate(r):
        col = RED if (i == 2 and float(v.strip('%')) < float(r[3].strip('%')) - 5) else NAVY
        tb(sl, v, xc[i], y, wc[i], 0.28, size=8, color=col)
    y += 0.31

# Charts right
img(sl, 'roc_curve_chart.png',      6.1, 2.7, 3.5, 2.2)
img(sl, 'score_band_chart.png',     6.1, 5.05, 7.0, 2.2)
img(sl, 'coef_vs_expert_chart.png', 9.75, 2.7, 3.5, 2.2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Q5: Model Evaluation
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'Q5: Model Evaluation Framework', 'Three pillars: goodness-of-fit, discriminatory power, calibration')

pillars = [
    (NAVY, 'Pillar 1\nGoodness-of-Fit', [
        ('Hosmer-Lemeshow Test', 'Stat=297.94, p≈0.00 → Poor fit\nCause: extreme WoE (±20) from sparse bins\nFix: Laplace smoothing + bin merging'),
        ('Binomial Test', 'Observed: 50 defaults (9.0%)\nPredicted avg PD: 34.97% → 4× overestimate\nFix: Platt scaling / PD recalibration'),
    ]),
    (TEAL, 'Pillar 2\nDiscriminatory Power', [
        ('AUC-ROC = 0.90', 'Model ranks 90% of default/non-default pairs correctly\nExcellent — robust even with extreme WoE values\nGini Coefficient = 2×AUC−1 = 0.80'),
        ('KS Statistic', 'Max separation between default/non-default CDFs\nScore means: 395.9 (default) vs 636.6 (non-default)\nStrong 240-point separation'),
    ]),
    (GOLD, 'Pillar 3\nForecasting Accuracy', [
        ('PSI — Population Stability', 'Monitor score distribution shift over time\nPSI < 0.1: stable; 0.1–0.25: minor shift; >0.25: rebuild'),
        ('Back-testing', 'Compare predicted PD vs realised default rates by band\nValidate annually; flag bands where |PD_pred − PD_actual| > 2%'),
    ]),
]
x = 0.4
for i, (color, title, items) in enumerate(pillars):
    add_rect(sl, x, 1.25, 4.1, 5.8, WHITE)
    add_rect(sl, x, 1.25, 4.1, 0.7, color)
    tb(sl, title, x+0.15, 1.3, 3.8, 0.6, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y = 2.05
    for metric, detail in items:
        add_rect(sl, x+0.15, y, 3.8, 0.28, LGRAY)
        tb(sl, metric, x+0.2, y+0.03, 3.7, 0.22, size=9, bold=True, color=NAVY)
        tb(sl, detail, x+0.2, y+0.33, 3.7, 0.7, size=8.5, color=NAVY)
        y += 1.2
    x += 4.3

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Score Distribution Deep Dive
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
page_bg(sl)
header(sl, 'Score Distribution Analysis', 'PDO-calibrated scores mapped to 300–850 range with band-level default rates')

img(sl, 'score_band_chart.png', 0.4, 1.25, 8.5, 4.5)

add_rect(sl, 9.1, 1.25, 3.85, 5.8, WHITE)
tb(sl, 'PDO Calibration', 9.3, 1.38, 3.5, 0.32, size=11, bold=True, color=NAVY)
divider(sl, 1.73)
pdo = [
    ('PDO', '20 (score doubles odds per 20 pts)'),
    ('Base Score', '600 at odds 50:1'),
    ('Factor', 'PDO/ln(2) = 28.85'),
    ('Offset', 'Base − Factor×ln(Base_Odds) = 487.1'),
    ('Intercept', 'Factor × β₀ = 10.55'),
    ('Score Range', '300–850 (min-max scaled)'),
]
y = 1.85
for k, v in pdo:
    tb(sl, k, 9.3, y, 1.4, 0.3, size=8.5, bold=True, color=TEAL)
    tb(sl, v, 10.75, y, 2.1, 0.3, size=8.5, color=NAVY)
    y += 0.38

divider(sl, 4.2)
tb(sl, 'Score Band Findings', 9.3, 4.32, 3.5, 0.3, size=10, bold=True, color=NAVY)
obs2 = [
    '300–498: 44.8% default rate → High Risk',
    '498–513: 5–13% default rate → Medium Risk',
    '513+: 0–3.5% default rate → Low Risk',
    'Clear separation validates model discrimination',
    'Bands 513–850: 0 defaults in top 5 bands',
]
y = 4.65
for o in obs2:
    tb(sl, f'• {o}', 9.3, y, 3.5, 0.36, size=8.5, color=NAVY)
    y += 0.38

add_rect(sl, 0.4, 5.9, 8.5, 0.85, WHITE)
tb(sl, '⚠  Calibration Issue: Predicted avg PD = 34.97% vs Actual 9.0% — due to class_weight="balanced" distorting probabilities. Apply Platt scaling before using PD values for provisioning or pricing.',
   0.6, 5.97, 8.1, 0.7, size=9, color=RED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 – Conclusion & Recommendations
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl, 0, 0, 0.35, 7.5, GOLD)

tb(sl, 'Conclusion & Recommendations', 1.0, 0.3, 12.0, 0.7, size=24, bold=True, color=WHITE)
divider(sl, 1.1)

recs_final = [
    (TEAL, 'Rebalance Feature Weights',
     'Increase C3 → 30%, D/E Ratio & DSCR → 15% each.\nRemove F3 Sale Growth (IV=0.04) — no predictive value.'),
    (TEAL, 'Fix WoE Calculation',
     'Apply Laplace smoothing (ε=0.5) to all bins.\nMerge bins with n<5 or zero-event cells.'),
    (TEAL, 'Recalibrate PD Probabilities',
     'class_weight="balanced" inflates PD 4×.\nUse Platt scaling or match base_odds to actual default rate.'),
    (GOLD, 'Validate Model Robustly',
     'HL test fails due to extreme WoE — fix binning first.\nImplement PSI monitoring for score drift detection.'),
    (GOLD, 'Add Macroeconomic Variables',
     'Integrate GDP growth, interest rate, oil price.\nTwo-stage: micro scorecard + macro adjustment factor.'),
    (GOLD, 'Rating Scale Calibration',
     'Define bands with log-linear PD spacing (PD doubles per grade).\nEnsure ≥20 defaults per band for stable PD estimates.'),
]
xs4 = [0.5, 4.65, 8.8]
ys4 = [1.3, 4.0]
idx = 0
for r in range(2):
    for c in range(3):
        color, title, detail = recs_final[idx]
        add_rect(sl, xs4[c], ys4[r], 3.95, 2.45, RGBColor(0x0A, 0x22, 0x44))
        add_rect(sl, xs4[c], ys4[r], 3.95, 0.06, color)
        tb(sl, title, xs4[c]+0.15, ys4[r]+0.12, 3.65, 0.35, size=10, bold=True, color=WHITE)
        tb(sl, detail, xs4[c]+0.15, ys4[r]+0.5, 3.65, 1.8, size=9, color=MGRAY)
        idx += 1

tb(sl, 'Model achieves AUC=0.90 — strong discrimination — but requires calibration fixes before production deployment.',
   0.5, 6.6, 12.4, 0.5, size=10, color=GOLD, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = '/Users/leekim/prj/deloitte/Credit_Scorecard_Presentation_v3.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
