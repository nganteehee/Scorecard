from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── THEME COLOURS ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1e, 0x3a, 0x5f)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xF0, 0xB3, 0x23)
LGRAY  = RGBColor(0xF4, 0xF6, 0xF9)
DKGRAY = RGBColor(0x44, 0x44, 0x44)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x1A, 0x7A, 0x4A)
BLUE   = RGBColor(0x2E, 0x86, 0xC1)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── HELPERS ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=NAVY, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txb

def add_para(tf, text, size=13, bold=False, color=DKGRAY,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p   = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return p

def slide_header(slide, title, subtitle=None):
    """Navy top bar with gold accent line"""
    add_rect(slide, 0, 0, 13.33, 1.2, NAVY)
    add_rect(slide, 0, 1.2, 13.33, 0.06, GOLD)
    add_text(slide, title, 0.4, 0.15, 12, 0.8,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.82, 12, 0.4,
                 size=13, color=GOLD, italic=True)

def bullet_box(slide, l, t, w, h, items, title=None,
               bg=LGRAY, title_bg=NAVY, item_size=12):
    """Rounded-ish box with optional title bar and bullet items"""
    add_rect(slide, l, t, w, h, fill=bg)
    yoff = 0
    if title:
        add_rect(slide, l, t, w, 0.38, fill=title_bg)
        add_text(slide, title, l+0.1, t+0.04, w-0.2, 0.34,
                 size=13, bold=True, color=WHITE)
        yoff = 0.42
    txb = slide.shapes.add_textbox(
        Inches(l+0.15), Inches(t+yoff),
        Inches(w-0.3),  Inches(h-yoff-0.1))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(item_size)
        run.font.color.rgb = DKGRAY

def kpi_box(slide, l, t, w, h, label, value, value_color=NAVY):
    add_rect(slide, l, t, w, h, fill=WHITE)
    # border
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.color.rgb = NAVY
    shape.line.width = Pt(1.5)
    add_text(slide, value, l, t+0.1, w, 0.55,
             size=30, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t+0.65, w, 0.35,
             size=10, color=DKGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl, 0, 4.8, 13.33, 0.08, GOLD)
add_text(sl, 'Credit Scorecard Development', 1, 1.5, 11, 1.2,
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, 'SME Portfolio — Quantitative Analysis & Model Enhancement',
         1, 2.8, 11, 0.6, size=20, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
add_text(sl, 'Candidate Case Study  |  Deloitte', 1, 3.6, 11, 0.5,
         size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, 'March 2026', 1, 5.3, 11, 0.5,
         size=13, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'Agenda')

agenda = [
    ('01', 'Background & Problem Statement'),
    ('02', 'Current State — Expert Weights'),
    ('03', 'Q1: Weight Assessment — Expert vs Data'),
    ('04', 'Q2: Scorecard Enhancement Suggestions'),
    ('05', 'Q3: EDA Approach'),
    ('06', 'Q4: Logistic Regression for PD Estimation'),
    ('07', 'Q5: Model Performance Evaluation'),
    ('08', 'Conclusion & Recommendations'),
]
for i, (num, label) in enumerate(agenda):
    row = i % 4
    col = i // 4
    l = 0.5 + col * 6.5
    t = 1.5 + row * 1.4
    add_rect(sl, l, t, 0.55, 0.8, NAVY)
    add_text(sl, num, l, t+0.05, 0.55, 0.7,
             size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_rect(sl, l+0.55, t, 5.6, 0.8, WHITE)
    add_text(sl, label, l+0.65, t+0.15, 5.4, 0.55,
             size=13, bold=False, color=NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BACKGROUND
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'Background & Problem Statement')

bullet_box(sl, 0.4, 1.4, 5.8, 2.6, title='Context',
    items=[
        '▸  SME loan portfolio of a bank',
        '▸  557 borrowers | 50 defaults (9% default rate)',
        '▸  8 features: 3 credit history + 5 financial factors',
        '▸  Hybrid scorecard: expert weights + statistical support',
    ])

bullet_box(sl, 6.6, 1.4, 6.3, 2.6, title='Problem Statement',
    items=[
        '▸  Current weights determined purely by expert judgement',
        '▸  Bank wants quantitative analysis to objectively support weights',
        '▸  Low default sample (50) → hybrid approach required',
        '▸  Goal: validate, refine, and enhance the scorecard',
    ])

bullet_box(sl, 0.4, 4.2, 12.5, 2.8, title='Current Expert Weights',
    items=[
        'C (Credit History) = 50%   →   C1 Debt Restructuring: 25%  |  C2 Debt Rescheduling: 25%  |  C3 Credit Bureau: 50%',
        'F (Financial Factor) = 50%  →   F1 Net Fixed Asset: 20%  |  F2 Net Profit Margin: 20%  |  F3 Sale Growth: 20%  |  F4 D/E Ratio: 20%  |  F5 DSCR: 20%',
        'Note: Weights of C+F = 100%, sub-weights within C and F each sum to 100%',
    ])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Q1: WEIGHT ASSESSMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'Q1: Weight Assessment — Expert vs Data-Driven',
             'Information Value (IV) used as objective measure of discriminatory power')

# Table header
cols = ['Feature', 'Expert Wt', 'IV', 'IV-based Wt', 'Gap', 'Verdict']
widths = [3.2, 1.4, 1.0, 1.6, 1.0, 2.0]
x = 0.4; y = 1.5
add_rect(sl, x, y, sum(widths)+0.1, 0.38, NAVY)
xx = x + 0.05
for c, w in zip(cols, widths):
    add_text(sl, c, xx, y+0.05, w, 0.3, size=11, bold=True, color=WHITE)
    xx += w

rows = [
    ('C1 Debt Restructuring',  '12.5%', '1.17', '8.3%',  '-4.2%', '⚠ Over-weighted',   WHITE,  RED),
    ('C2 Debt Rescheduling',   '12.5%', '0.57', '4.0%',  '-8.5%', '⚠ Over-weighted',   WHITE,  RED),
    ('C3 Credit Bureau',       '25.0%', '5.43', '38.4%', '+13.4%','✓ Under-weighted',  WHITE,  GREEN),
    ('F1 Net Fixed Asset',     '10.0%', '2.60', '18.4%', '+8.4%', '✓ Under-weighted',  WHITE,  GREEN),
    ('F2 Net Profit Margin',   '10.0%', '1.28', '9.1%',  '-0.9%', '~ Correct',         WHITE,  BLUE),
    ('F3 Sale Growth',         '10.0%', '0.04', '0.3%',  '-9.7%', '✗ Remove',          WHITE,  RED),
    ('F4 Debt-to-Equity',      '10.0%', '1.45', '10.4%', '+0.4%', '~ Correct',         WHITE,  BLUE),
    ('F5 DSCR',                '10.0%', '1.58', '11.2%', '+1.2%', '~ Correct',         WHITE,  BLUE),
]

for i, row in enumerate(rows):
    bg = WHITE if i % 2 == 0 else LGRAY
    add_rect(sl, x, y+0.38+i*0.42, sum(widths)+0.1, 0.42, bg)
    xx = x + 0.05
    for j, (val, w) in enumerate(zip(row[:6], widths)):
        col_c = row[7] if j == 5 else DKGRAY
        bold  = j == 5
        add_text(sl, val, xx, y+0.4+i*0.42, w, 0.38,
                 size=10, color=col_c, bold=bold)
        xx += w

add_text(sl, '* IV > 0.3 = Strong predictor  |  IV < 0.1 = Weak/useless',
         0.4, 7.1, 12, 0.3, size=9, color=DKGRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Q1 SUPPORTING EVIDENCE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'Q1: Key Evidence — Default Rates by Feature',
             'Data confirms C3 dominates; F3 Sale Growth shows near-zero discriminatory power')

kpi_box(sl, 0.4,  1.5, 2.8, 1.0, 'C3: Overdue >90d', '100%', RED)
kpi_box(sl, 3.4,  1.5, 2.8, 1.0, 'C1: Restructured <6m', '62.5%', RED)
kpi_box(sl, 6.4,  1.5, 2.8, 1.0, 'F3: Sale Growth >30%', '11.4%', DKGRAY)
kpi_box(sl, 9.4,  1.5, 3.1, 1.0, 'F3: Sale Growth <=0', '9.1%', DKGRAY)

bullet_box(sl, 0.4, 2.75, 5.9, 2.1, title='C3 — Severely Under-weighted (IV=5.43)',
    items=[
        '▸  "Overdue >90d >1yr" → 100% default rate — 7/7 customers defaulted',
        '▸  "No delinquency" → 5.6% default — clear separation',
        '▸  Chi-square p < 0.001 — strongest statistical association',
        '▸  IV-based weight: 38.4% vs current 25% → gap of +13.4%',
    ], item_size=11)

bullet_box(sl, 6.6, 2.75, 6.3, 2.1, title='F3 Sale Growth — Should be Removed (IV=0.04)',
    items=[
        '▸  Default rates flat across all bins: 6.8% to 12.5% — no clear pattern',
        '▸  Chi-square p = 0.758 — not statistically significant',
        '▸  Spearman ρ = +0.696 — directionally WRONG (higher growth = more default)',
        '▸  IV = 0.04 → below 0.02 threshold for useful predictors',
    ], item_size=11)

bullet_box(sl, 0.4, 5.05, 12.5, 1.9, title='Recommended Weight Revision',
    items=[
        'C (Credit History): 55%   →   C1: 20%  |  C2: 10%  |  C3: 70%     (C3 receives majority weight given IV=5.43)',
        'F (Financial Factor): 45%  →   F1: 30%  |  F2: 15%  |  F3: 0% (remove)  |  F4: 25%  |  F5: 30%',
    ], item_size=11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Q2: SCORECARD ENHANCEMENTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'Q2: Scorecard Enhancement Suggestions',
             'Six data-driven issues identified through monotonicity analysis and WoE calibration')

issues = [
    ('F3 Sale Growth\nInverted Direction',
     'Higher growth → higher default (11.4% vs 6.9%)\nSpearman ρ = +0.696 — scoring direction wrong',
     'Remove from scorecard or assign flat/neutral scores'),
    ('F4 Debt-to-Equity\nWrong Peak Bin',
     '(0, 0.5) gets max +100 pts but has 20% default rate\n(0.5, 1) has 0% defaults but only +80 pts',
     'Shift max score to (0.5, 1); redesign bin boundaries'),
    ('F2 Net Profit Margin\nNon-Monotone',
     '>15% margin defaults at 12% — same as loss-making firms\nNot statistically significant (p=0.176)',
     'Cap max score at (7.5–15) bin; reduce score for >15%'),
    ('C3 Bureau: Ambiguous\n"Not Performed" Rule',
     '"Not performed (exempted)" = +100 pts, 0% default\n"Not performed" = -100 pts, 16.7% default\n200-point gap for visually identical categories',
     'Document clear business rule distinguishing the two'),
    ('F5 DSCR: Missing\nData Scored as Neutral',
     '122 customers with missing DSCR default at 14.8%\nNearly as risky as worst bin — yet scored 0 (neutral)',
     'Assign penalty score similar to (0, 0.25) bin'),
    ('Replace Expert Scores\nwith PDO Calibration',
     'Current round numbers (-100, 0, +75, +100) are arbitrary\nPDO framework: Points = -(Factor × WoE)\nExample: C3 "Overdue >90d" should be -608 not -100',
     'Adopt PDO-based scores: Base=600, PDO=20, Odds=50:1'),
]

for i, (title, problem, fix) in enumerate(issues):
    col = i % 2
    row = i // 2
    l = 0.4 + col * 6.5
    t = 1.45 + row * 1.95
    add_rect(sl, l, t, 6.2, 1.85, WHITE)
    add_rect(sl, l, t, 6.2, 0.38, NAVY)
    add_text(sl, title, l+0.1, t+0.04, 6.0, 0.34, size=11, bold=True, color=WHITE)
    add_text(sl, '⚠ ' + problem, l+0.1, t+0.42, 6.0, 0.8, size=9, color=RED)
    add_text(sl, '→ ' + fix,     l+0.1, t+1.22, 6.0, 0.55, size=9, color=GREEN, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Q3: EDA APPROACH
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'Q3: EDA Approach for Default Prediction',
             'Structured four-step framework to identify key predictors')

steps = [
    ('(1) Data Understanding\n& Quality Check',
     ['Variable types: categorical, numerical, macro time-series',
      'Class balance: 50 defaults / 557 (9%) — significant imbalance',
      'Missing vs sentinel values (99999998/99999999) → treat as bins',
      'Outliers & skewness: IQR-based detection; cap at 95th pctile']),
    ('(2) Univariate Analysis',
     ['Categorical: default rate per bin → concentration check',
      'Numerical: boxplots split by Default Flag → visual separation',
      'Monotonicity: default rate should move consistently per bin',
      'Macro variables: time-series plot vs portfolio default rate']),
    ('(3) Bivariate Analysis',
     ['Numerical vs default: KS statistic, point-biserial correlation',
      'Categorical vs default: chi-square test (F2 p=0.18, F3 p=0.76 → NS)',
      'Multicollinearity: Pearson matrix + Cramér\'s V for categoricals',
      'Vintage analysis: default rate by origination year → macro effects']),
    ('(4) WoE & IV Analysis',
     ['WoE_i = ln(%Bad_i / %Good_i)   |   IV = Σ(%Bad - %Good) × WoE',
      'IV ranking: C3=5.43 ▸ F1=2.60 ▸ F5=1.58 ▸ F4=1.45 ▸ F3=0.04',
      'Features with IV < 0.02 excluded; non-monotone WoE → rebin',
      'Objective basis for weight assignment — replaces expert judgement']),
]

for i, (title, bullets) in enumerate(steps):
    l = 0.4 + (i % 2) * 6.5
    t = 1.45 + (i // 2) * 2.8
    add_rect(sl, l, t, 6.2, 2.6, WHITE)
    add_rect(sl, l, t, 6.2, 0.45, NAVY)
    add_text(sl, f'Step {i+1}  |  {title}', l+0.1, t+0.05, 6.0, 0.4,
             size=12, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        add_text(sl, '▸  ' + b, l+0.1, t+0.5+j*0.5, 6.0, 0.48,
                 size=10, color=DKGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Q4: LOGISTIC REGRESSION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'Q4: Logistic Regression for PD Estimation',
             'Six-step development process with WoE encoding and PDO calibration')

steps_lr = [
    ('1. Data Prep',    'Handle sentinels as separate bins\nFill nulls → "Not available" bin'),
    ('2. WoE Encode',   'Replace bins with log-odds values\nLinearises LR assumption by construction'),
    ('3. Train/Test',   'Stratified 80/20 split\n5-fold CV — only 50 defaults available'),
    ('4. Train Model',  'class_weight="balanced" for 9% imbalance\nL2 regularisation — tune C via grid search'),
    ('5. PDO Calibrate','Points = -(Factor × β × WoE)\nBase=600, PDO=20, Odds=50:1'),
    ('6. Score → PD',   'PD = 1 / (1 + e^((score-offset)/factor))\nInterpretable probability output'),
]

for i, (title, body) in enumerate(steps_lr):
    l = 0.4 + (i % 3) * 4.25
    t = 1.45 + (i // 3) * 2.0
    add_rect(sl, l, t, 4.0, 1.8, WHITE)
    add_rect(sl, l, t, 4.0, 0.42, NAVY)
    add_text(sl, title, l+0.1, t+0.06, 3.8, 0.36, size=13, bold=True, color=WHITE)
    add_text(sl, body, l+0.1, t+0.5, 3.8, 1.2, size=11, color=DKGRAY)

# Assumptions box
bullet_box(sl, 0.4, 5.6, 12.5, 1.7, title='Key Assumptions & Limitations',
    items=[
        'Assumptions: Binary outcome (Bernoulli) | Linearity in log-odds (satisfied via WoE) | Independence of observations | No perfect multicollinearity | Correct model specification',
        'Key limitation: 50 defaults / 8 features = 6.25 EPV (below recommended threshold of 10) → use stronger regularisation, reduce features, report confidence intervals',
    ], item_size=10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MODEL RESULTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'Model Results — Logistic Regression (WoE-encoded)',
             'AUC = 0.87 | Gini = 0.74 | 5-Fold CV AUC = 0.82 ± 0.07')

kpi_box(sl, 0.4,  1.5, 2.8, 1.1, 'AUC-ROC', '0.87', GREEN)
kpi_box(sl, 3.5,  1.5, 2.8, 1.1, 'Gini Coefficient', '0.74', GREEN)
kpi_box(sl, 6.6,  1.5, 2.8, 1.1, 'CV AUC (5-fold)', '0.82', BLUE)
kpi_box(sl, 9.7,  1.5, 3.0, 1.1, 'Good Mean Score', '785 pts', NAVY)

# LR weights table
add_rect(sl, 0.4, 2.85, 5.8, 0.38, NAVY)
for j, h in enumerate(['Feature', 'LR Coef', 'LR Weight']):
    add_text(sl, h, 0.5 + j*1.9, 2.88, 1.8, 0.32, size=11, bold=True, color=WHITE)

coef_rows = [
    ('C3 Credit Bureau',     '0.890', '24.4%', GREEN),
    ('F4 Debt-to-Equity',    '0.754', '20.7%', GREEN),
    ('F5 DSCR',              '0.583', '16.0%', GREEN),
    ('C2 Rescheduling',      '0.312', ' 8.6%', DKGRAY),
    ('C1 Restructuring',     '0.308', ' 8.4%', DKGRAY),
    ('F3 Sale Growth',       '0.297', ' 8.1%', RED),
    ('F1 Net Fixed Asset',   '0.285', ' 7.8%', DKGRAY),
    ('F2 Net Profit Margin', '0.220', ' 6.0%', DKGRAY),
]
for i, (feat, coef, wt, col) in enumerate(coef_rows):
    bg = WHITE if i%2==0 else LGRAY
    add_rect(sl, 0.4, 3.23+i*0.42, 5.8, 0.42, bg)
    add_text(sl, feat, 0.5,  3.25+i*0.42, 1.85, 0.38, size=10, color=DKGRAY)
    add_text(sl, coef, 2.4,  3.25+i*0.42, 1.85, 0.38, size=10, color=DKGRAY)
    add_text(sl, wt,   4.3,  3.25+i*0.42, 1.85, 0.38, size=10, color=col, bold=True)

bullet_box(sl, 6.6, 2.85, 6.3, 4.5, title='Score Distribution & Interpretation',
    items=[
        '▸  Good customers (Default=0): Mean score = 785 pts',
        '▸  Bad customers  (Default=1): Mean score = 344 pts',
        '▸  Separation gap: ~440 points — strong discrimination',
        '',
        '▸  Score → PD mapping (Base=600, PDO=20, Odds=50:1):',
        '       Score 645  →  PD ≈ 1.0%   (very low risk)',
        '       Score 600  →  PD ≈ 2.0%   (base odds)',
        '       Score 572  →  PD ≈ 5.0%   (moderate risk)',
        '       Score 511  →  PD ≈ 20.0%  (high risk)',
        '',
        '▸  LR coefficients are directional — not true linear weights',
        '▸  IV remains the primary basis for scorecard weight justification',
    ], item_size=10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Q5: MODEL EVALUATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'Q5: Model Performance Evaluation',
             'Three pillars: Discriminatory Power | Goodness-of-Fit | Forecasting Accuracy')

pillar_data = [
    ('Discriminatory Power', NAVY, [
        'AUC-ROC: area under the ROC curve',
        '→ AUC = 0.87 (target > 0.75) ✓',
        'Gini = 2×AUC - 1 = 0.74',
        'KS Statistic: max separation between',
        '   good and bad score distributions',
        '→ Target KS > 0.30',
    ]),
    ('Goodness of Fit', BLUE, [
        'Log-likelihood: measures how well',
        '   model fits observed defaults',
        'AIC / BIC: penalises model complexity',
        '   → use to compare lag structures',
        'Hosmer-Lemeshow test: checks if',
        '   predicted PD = actual default rate',
    ]),
    ('Forecasting Accuracy', GREEN, [
        'Confusion matrix: TP, FP, FN, TN',
        'Sensitivity (recall): % defaults caught',
        'Precision: % flagged that truly default',
        'PSI < 0.10: score distribution stable',
        'Calibration curve: predicted vs actual PD',
        'Backtesting: compare PD to realised DR',
    ]),
]
for i, (title, color, items) in enumerate(pillar_data):
    l = 0.4 + i * 4.3
    add_rect(sl, l, 1.45, 4.0, 0.45, color)
    add_text(sl, title, l+0.1, 1.48, 3.8, 0.4, size=13, bold=True, color=WHITE)
    add_rect(sl, l, 1.9, 4.0, 3.8, WHITE)
    for j, item in enumerate(items):
        bold = item.startswith('→')
        col  = GREEN if item.startswith('→') else DKGRAY
        add_text(sl, item, l+0.15, 1.95+j*0.6, 3.7, 0.55,
                 size=10, color=col, bold=bold)

bullet_box(sl, 0.4, 5.85, 12.5, 1.45, title='Validation Framework for Small Sample (50 Defaults)',
    items=[
        '▸  Use 5-fold stratified CV instead of single train/test split — maximises use of limited default data',
        '▸  Bootstrap confidence intervals for AUC — quantify uncertainty in performance estimates',
        '▸  Report sensitivity ≥ 80% as minimum threshold — missing defaults is more costly than false alarms in credit risk',
    ], item_size=10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SCORECARD BIN SCORES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
slide_header(sl, 'PDO-Calibrated Scorecard — Selected Bins',
             'Points = -(Factor × β × WoE)  |  Base Score=600, PDO=20, Odds=50:1')

add_text(sl, 'C3: Credit Bureau (Most Impactful — IV=5.43)',
         0.4, 1.45, 6.0, 0.35, size=12, bold=True, color=NAVY)
c3_rows = [
    ('Overdue >90d >1yr ago',       '100%',  '+21.06', '-608', '-100'),
    ('Overdue 61-90d <1yr',         '30.0%', ' +1.47', ' -42', ' -75'),
    ('Overdue 31-60d <1yr',         '17.2%', ' +0.75', ' -22', ' -50'),
    ('No delinquency history',      ' 5.6%', ' -0.52', ' +15', '+100'),
    ('No record found',             ' 0.0%', '-19.63', '+566', '   0'),
]
hdrs = ['Bin', 'Default Rate', 'WoE', 'PDO Pts', 'Expert Pts']
add_rect(sl, 0.4, 1.85, 6.2, 0.35, NAVY)
for j, (h, w) in enumerate(zip(hdrs, [2.4, 1.2, 0.9, 0.9, 1.0])):
    xx = 0.5 + sum([2.4,1.2,0.9,0.9,1.0][:j])
    add_text(sl, h, xx, 1.87, w, 0.3, size=9, bold=True, color=WHITE)

for i, row in enumerate(c3_rows):
    bg = WHITE if i%2==0 else LGRAY
    add_rect(sl, 0.4, 2.2+i*0.42, 6.2, 0.42, bg)
    xx = 0.5
    for j, (val, w) in enumerate(zip(row, [2.4,1.2,0.9,0.9,1.0])):
        col = RED if (j in [2,3] and val.strip().startswith('+') and j==2) else DKGRAY
        if j == 3:
            col = RED if float(val) < 0 else GREEN
        elif j == 4:
            col = DKGRAY
        add_text(sl, val, xx, 2.22+i*0.42, w, 0.38, size=9, color=col)
        xx += w

add_text(sl, 'F4: Debt-to-Equity (Scoring Direction Issue)',
         6.8, 1.45, 6.0, 0.35, size=12, bold=True, color=NAVY)
f4_rows = [
    ('(3, 5)',    '33.3%', '+1.62', '-47', '+20'),
    ('(0, 0.5)', '20.0%', '+0.93', '-27', '+100 ← WRONG'),
    ('(1.5, 3)', '20.0%', '+0.93', '-27', '+40'),
    ('(0.5, 1)', ' 0.0%', '-20.23','+584','+80'),
    ('(1, 1.5)', ' 3.8%', '-0.90', '+26', '+60'),
]
add_rect(sl, 6.8, 1.85, 6.1, 0.35, NAVY)
for j, (h, w) in enumerate(zip(hdrs, [1.8,1.2,1.0,0.9,1.2])):
    xx = 6.9 + sum([1.8,1.2,1.0,0.9,1.2][:j])
    add_text(sl, h, xx, 1.87, w, 0.3, size=9, bold=True, color=WHITE)
for i, row in enumerate(f4_rows):
    bg = WHITE if i%2==0 else LGRAY
    add_rect(sl, 6.8, 2.2+i*0.42, 6.1, 0.42, bg)
    xx = 6.9
    for j, (val, w) in enumerate(zip(row, [1.8,1.2,1.0,0.9,1.2])):
        col = RED if 'WRONG' in val else DKGRAY
        if j == 3:
            col = RED if val.strip().startswith('-') else GREEN
        add_text(sl, val, xx, 2.22+i*0.42, w, 0.38, size=9, color=col, bold='WRONG' in val)
        xx += w

bullet_box(sl, 0.4, 4.45, 12.5, 1.5, title='Key Insight: PDO vs Expert Scores',
    items=[
        '▸  C3 "No record found" = 0 expert pts  vs  +566 PDO pts — expert scoring leaves significant signal uncaptured',
        '▸  F4 "(0, 0.5)" = +100 expert pts  vs  -27 PDO pts — expert scoring penalises safest bin incorrectly',
        '▸  PDO calibration replaces arbitrary round numbers with mathematically grounded, auditable scores — critical for regulatory review',
    ], item_size=10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl, 0, 5.5, 13.33, 0.08, GOLD)
add_text(sl, 'Conclusion & Recommendations', 0.5, 0.15, 12, 0.8,
         size=28, bold=True, color=WHITE)
add_text(sl, 'Data-driven enhancements to move beyond pure expert judgement',
         0.5, 0.82, 12, 0.4, size=13, color=GOLD, italic=True)

recs = [
    ('Reweight C3',       'Increase Credit Bureau weight from 25% to ~38%\nbased on IV=5.43 — dominant predictor'),
    ('Remove F3',         'Sale Growth has IV=0.04, p=0.758, inverted direction\n→ remove or assign flat score'),
    ('Fix F4 Bins',       'Restructure D/E bins — (0.5,1) has 0% default\nbut currently underscored vs (0,0.5)'),
    ('PDO Calibration',   'Replace round expert numbers with WoE-based\nPDO scores — interpretable and auditable'),
    ('Address EPV',       '50 defaults / 8 features = 6.25 EPV < 10\n→ reduce features, increase regularisation'),
    ('Add Macro Layer',   'Dataset lacks time dimension — recommend adding\nGDP growth (lag 4Q) and interest rate overlay'),
]

for i, (title, body) in enumerate(recs):
    col = i % 3
    row = i // 3
    l = 0.4 + col * 4.3
    t = 1.4 + row * 1.85
    add_rect(sl, l, t, 4.0, 1.7,  RGBColor(0x1a, 0x2e, 0x4a))
    add_rect(sl, l, t, 4.0, 0.42, GOLD)
    add_text(sl, title, l+0.1, t+0.05, 3.8, 0.35, size=12, bold=True, color=NAVY)
    add_text(sl, body,  l+0.1, t+0.5,  3.8, 1.1,  size=10, color=WHITE)

add_text(sl, 'The hybrid approach — combining quantitative IV/WoE analysis with expert judgement — is the appropriate framework given the small sample of 50 defaults.',
         0.5, 5.7, 12, 0.6, size=11, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────────
path = '/Users/leekim/prj/deloitte/Credit_Scorecard_Presentation.pptx'
prs.save(path)
print(f'Saved: {path}')
