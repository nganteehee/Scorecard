"""
Credit Scorecard Presentation — Professional Version
12 slides, logical flow, clean design
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── PALETTE ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x2B, 0x55)   # deep navy
TEAL    = RGBColor(0x00, 0x7A, 0x8A)   # accent teal
GOLD    = RGBColor(0xE8, 0xA0, 0x00)   # gold
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT  = RGBColor(0xF7, 0xF9, 0xFC)
LGRAY   = RGBColor(0xE8, 0xEC, 0xF1)
MGRAY   = RGBColor(0x8A, 0x96, 0xA8)
DKGRAY  = RGBColor(0x2D, 0x3A, 0x4A)
RED     = RGBColor(0xC0, 0x39, 0x2B)
LRED    = RGBColor(0xFD, 0xED, 0xEC)
GREEN   = RGBColor(0x1A, 0x73, 0x48)
LGREEN  = RGBColor(0xE9, 0xF7, 0xEF)
AMBER   = RGBColor(0xF3, 0x9C, 0x12)
LAMBER  = RGBColor(0xFE, 0xF9, 0xE7)
BLUE    = RGBColor(0x1A, 0x5C, 0x9A)
LBLUE   = RGBColor(0xEB, 0xF5, 0xFB)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── PRIMITIVE HELPERS ─────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=NAVY, border=None, border_w=1):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width     = Pt(border_w)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       sz=12, bold=False, italic=False,
       color=DKGRAY, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def add_run(tf, text, sz=11, bold=False, italic=False,
            color=DKGRAY, align=PP_ALIGN.LEFT, space=3):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space)
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

# ── COMPOUND HELPERS ──────────────────────────────────────────────────────────
def page_bg(slide, color=OFFWHT):
    rect(slide, 0, 0, 13.33, 7.5, fill=color)

def header_bar(slide, title, subtitle=None, accent=TEAL):
    """Left navy sidebar + top accent line"""
    rect(slide, 0, 0, 13.33, 1.05, NAVY)
    rect(slide, 0, 1.05, 13.33, 0.05, accent)
    tb(slide, title, 0.45, 0.08, 12.4, 0.65,
       sz=24, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.45, 0.68, 12.4, 0.35,
           sz=11, italic=True, color=GOLD)

def section_label(slide, text, l, t, w=4, color=TEAL):
    tb(slide, text.upper(), l, t, w, 0.28,
       sz=8, bold=True, color=color)

def card(slide, l, t, w, h,
         title=None, title_color=NAVY, title_bg=None,
         bg=WHITE, border_color=LGRAY):
    rect(slide, l, t, w, h, fill=bg, border=border_color, border_w=0.5)
    if title:
        title_bg = title_bg or NAVY
        rect(slide, l, t, w, 0.36, fill=title_bg)
        tb(slide, title, l+0.12, t+0.05, w-0.24, 0.28,
           sz=11, bold=True, color=WHITE)

def bullet_card(slide, l, t, w, h, title, items,
                title_bg=NAVY, bg=WHITE, sz=10, gap=0.44):
    card(slide, l, t, w, h, title=title, title_bg=title_bg, bg=bg)
    for i, item in enumerate(items):
        tb(slide, item, l+0.14, t+0.42+i*gap, w-0.28, gap-0.04,
           sz=sz, color=DKGRAY)

def kpi(slide, l, t, w, h, value, label,
        val_color=NAVY, val_sz=32, bg=WHITE):
    rect(slide, l, t, w, h, fill=bg, border=LGRAY, border_w=0.5)
    tb(slide, value, l, t+0.1,  w, h*0.58,
       sz=val_sz, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    tb(slide, label, l, t+h*0.62, w, h*0.35,
       sz=9, color=MGRAY, align=PP_ALIGN.CENTER)

def divider(slide, l, t, w=12.4, color=LGRAY, h=0.02):
    rect(slide, l, t, w, h, fill=color)

def tag(slide, text, l, t, bg=TEAL, color=WHITE, sz=9):
    w = len(text)*0.085 + 0.25
    rect(slide, l, t, w, 0.26, fill=bg)
    tb(slide, text, l+0.08, t+0.04, w-0.1, 0.2,
       sz=sz, bold=True, color=color)
    return w


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 01 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl, NAVY)
rect(sl, 0, 0,    0.5,  7.5, TEAL)
rect(sl, 0, 5.85, 13.33, 0.06, GOLD)
tb(sl, 'CREDIT SCORECARD', 1.0, 1.6, 11, 0.9,
   sz=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
tb(sl, 'DEVELOPMENT', 1.0, 2.45, 11, 0.9,
   sz=44, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
rect(sl, 1.0, 3.45, 5.5, 0.04, TEAL)
tb(sl, 'SME Portfolio  —  Quantitative Analysis & Model Enhancement',
   1.0, 3.6, 11, 0.5, sz=15, italic=True, color=WHITE)
tb(sl, 'Candidate Case Study  |  Deloitte  |  March 2026',
   1.0, 6.1, 11, 0.4, sz=11, color=GOLD)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 02 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'Executive Summary',
           'Key findings and recommendations at a glance')

# KPIs
kpi(sl, 0.45, 1.2, 2.7, 1.05, '557',  'Total Borrowers',       NAVY)
kpi(sl, 3.3,  1.2, 2.7, 1.05, '50',   'Observed Defaults',     RED)
kpi(sl, 6.15, 1.2, 2.7, 1.05, '8.98%','Portfolio Default Rate', AMBER)
kpi(sl, 9.0,  1.2, 2.7, 1.05, '0.87', 'Model AUC-ROC',         GREEN)

divider(sl, 0.45, 2.38)

# Three finding columns
cols = [
    ('Weight Misalignment', TEAL, [
        '▸ C3 Credit Bureau under-weighted by 13.4%',
        '▸ F3 Sale Growth over-weighted by 9.7% — useless predictor (IV=0.04)',
        '▸ IV analysis provides objective basis for reweighting',
    ]),
    ('Scorecard Issues', AMBER, [
        '▸ F4 Debt-to-Equity scoring direction is inverted',
        '▸ F3 Sale Growth shows non-monotone, inverted default pattern',
        '▸ PDO calibration replaces arbitrary round-number scores',
    ]),
    ('Model Performance', GREEN, [
        '▸ Logistic regression on WoE features: AUC = 0.87',
        '▸ Good vs Bad score gap: 785 vs 344 pts (441 pt separation)',
        '▸ EPV = 6.25 < 10 → hybrid expert + statistical approach appropriate',
    ]),
]
for i, (title, color, items) in enumerate(cols):
    l = 0.45 + i*4.3
    rect(sl, l, 2.55, 4.0, 0.06, color)
    tb(sl, title, l, 2.7, 4.0, 0.35, sz=12, bold=True, color=DKGRAY)
    for j, item in enumerate(items):
        tb(sl, item, l, 3.1+j*0.52, 4.0, 0.5, sz=10, color=DKGRAY)

divider(sl, 0.45, 5.55)
tb(sl, 'Bottom line: The hybrid approach — anchoring on expert judgement while using IV/WoE analysis to validate and refine weights — is the most appropriate framework given the limited default sample of 50 observations.',
   0.45, 5.65, 12.4, 0.6, sz=10, italic=True, color=MGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 03 — BACKGROUND & DATA OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'Background & Data Overview',
           'SME portfolio scorecard — hybrid expert-statistical development approach')

# Left: context
bullet_card(sl, 0.45, 1.2, 5.9, 2.05, 'Business Context', [
    '▸  SME loan portfolio of a regional bank',
    '▸  Low default sample (50 cases) → pure statistical approach unreliable',
    '▸  Hybrid approach: expert weights anchored + quantitative refinement',
    '▸  Goal: objectively justify and improve expert-determined weights',
], sz=11, gap=0.42)

# Right: data overview
bullet_card(sl, 6.6, 1.2, 6.3, 2.05, 'Dataset Structure', [
    '▸  557 borrowers  |  50 defaults (9.0%)  |  507 non-defaults',
    '▸  3 credit history features (C1, C2, C3) — categorical',
    '▸  5 financial ratio features (F1–F5) — numerical, pre-binned',
    '▸  Special values: 99999998 (neg. equity), 99999999 (zero denominator)',
], sz=11, gap=0.42)

# Feature table
add_lbl = [
    ('C1', 'Debt Restructuring History',  'Categorical', '7 bins'),
    ('C2', 'Debt Rescheduling History',   'Categorical', '4 bins + NaN'),
    ('C3', 'Credit Bureau Delinquency',   'Categorical', '12 bins'),
    ('F1', 'Net Fixed Asset',             'Numerical',   '5 bins'),
    ('F2', 'Net Profit Margin (%)',       'Numerical',   '6 bins + sentinel'),
    ('F3', 'Sale Growth (%)',             'Numerical',   '5 bins + sentinel'),
    ('F4', 'Debt-to-Equity Ratio',        'Numerical',   '8 bins + sentinel'),
    ('F5', 'Debt Service Coverage Ratio', 'Numerical',   '8 bins'),
]
rect(sl, 0.45, 3.42, 12.45, 0.36, NAVY)
for j, (h, w) in enumerate(zip(['ID','Feature Name','Type','Bins'],
                                 [0.5, 5.5, 2.8, 2.8])):
    xx = 0.55 + sum([0.5,5.5,2.8,2.8][:j])
    tb(sl, h, xx, 3.45, w, 0.3, sz=10, bold=True, color=WHITE)

for i, row in enumerate(add_lbl):
    bg = WHITE if i%2==0 else OFFWHT
    rect(sl, 0.45, 3.78+i*0.42, 12.45, 0.42, fill=bg, border=LGRAY, border_w=0.3)
    ws = [0.5, 5.5, 2.8, 2.8]
    xx = 0.55
    for j, (val, w) in enumerate(zip(row, ws)):
        col = TEAL if j==0 else DKGRAY
        bold = j==0
        tb(sl, val, xx, 3.8+i*0.42, w, 0.38, sz=10, color=col, bold=bold)
        xx += w

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 04 — EDA: KEY INSIGHTS (Q3 APPLIED)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'EDA: Key Insights from the Data',
           'Q3 applied — four-step framework reveals critical patterns')

# Step labels
steps = ['(1) Data Quality', '(2) Univariate', '(3) Bivariate', '(4) WoE & IV']
for i, s in enumerate(steps):
    l = 0.45 + i*3.25
    rect(sl, l, 1.18, 3.05, 0.36, TEAL if i==3 else LGRAY)
    col = WHITE if i==3 else MGRAY
    tb(sl, f'Step {i+1}', l+0.1, 1.2, 1.0, 0.3, sz=8, bold=True, color=col)
    tb(sl, s, l+0.7, 1.2, 2.3, 0.3, sz=9, bold=(i==3), color=col)

# Four finding boxes
findings = [
    ('Data Quality Check', TEAL, [
        '▸  8.98% overall default rate — class imbalance',
        '▸  282 nulls in C2 (51%) → treated as "Not available" bin',
        '▸  Sentinel 99999998/99999999 → separate risk bins',
        '▸  No duplicate customer records found',
        '▸  All 5 numerical features are right-skewed',
    ]),
    ('Univariate Findings', TEAL, [
        '▸  C1 "Restructured <6m" → 62.5% default rate',
        '▸  C3 "Overdue >90d" → 100% default rate (7/7)',
        '▸  F5 DSCR (0, 0.25) → 18.1% default — highest of financials',
        '▸  F3 Sale Growth: flat default rates across ALL bins',
        '▸  Smaller firms (NFA <100M) default at 12–15% vs 0–4% larger',
    ]),
    ('Bivariate Findings', TEAL, [
        '▸  C1, C2, C3 chi-square: p < 0.001 — highly significant',
        '▸  F2 Net Profit Margin: p = 0.176 — NOT significant',
        '▸  F3 Sale Growth: p = 0.758 — NOT significant',
        '▸  C1 & C2 likely correlated (Cramér\'s V) — co-occurrence',
        '▸  F4 D/E: non-monotone — (0,0.5) riskier than (0.5,1)',
    ]),
    ('WoE & IV Results', NAVY, [
        '▸  C3: IV = 5.43 → Very Strong (dominant predictor)',
        '▸  F1: IV = 2.63 → Strong (firm size matters)',
        '▸  F5: IV = 1.66 → Strong (debt coverage signal)',
        '▸  F4: IV = 1.48 → Strong (non-monotone — rebin needed)',
        '▸  F3: IV = 0.62 → Misleading (inverted direction)',
    ]),
]
for i, (title, tbg, items) in enumerate(findings):
    l = 0.45 + (i%2)*6.45
    t = 1.7  + (i//2)*2.75
    bullet_card(sl, l, t, 6.2, 2.55, title, items,
                title_bg=tbg, sz=10, gap=0.42)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 05 — Q1: WEIGHT ASSESSMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'Q1: Feature Weight Assessment',
           'Methodology: Information Value (IV) as objective measure of discriminatory power')

# Methodology note
rect(sl, 0.45, 1.2, 12.45, 0.54, OFFWHT, border=LGRAY, border_w=0.5)
tb(sl, 'IV = Σ (%Bad_i − %Good_i) × WoE_i    where    WoE_i = ln(%Bad_i / %Good_i)',
   0.6, 1.25, 7.0, 0.3, sz=11, bold=True, color=NAVY)
tb(sl, 'IV < 0.10 = Weak   |   0.10–0.30 = Medium   |   0.30–0.50 = Strong   |   > 0.50 = Very Strong',
   0.6, 1.52, 11.8, 0.2, sz=9, italic=True, color=MGRAY)

# Table
hdrs = ['Feature', 'Expert Weight', 'IV Score', 'IV-based Weight', 'Gap', 'Verdict']
ws   = [3.1, 1.6, 1.2, 1.9, 1.1, 2.25]
rect(sl, 0.45, 1.88, sum(ws)+0.3, 0.36, NAVY)
xx = 0.55
for h, w in zip(hdrs, ws):
    tb(sl, h, xx, 1.9, w, 0.3, sz=10, bold=True, color=WHITE)
    xx += w

rows = [
    ('C1  Debt Restructuring History',       '12.5%','1.17','8.3%', '−4.2%','Over-weighted',  RED,   LRED),
    ('C2  Debt Rescheduling History',         '12.5%','0.59','4.2%', '−8.3%','Over-weighted',  RED,   LRED),
    ('C3  Credit Bureau Delinquency',         '25.0%','5.43','38.6%','+13.6%','Under-weighted', GREEN, LGREEN),
    ('F1  Net Fixed Asset',                   '10.0%','2.63','18.7%','+ 8.7%','Under-weighted', GREEN, LGREEN),
    ('F2  Net Profit Margin (%)',             '10.0%','1.70','12.1%','+ 2.1%','Broadly correct', TEAL, OFFWHT),
    ('F3  Sale Growth (%)',                   '10.0%','0.62','4.4%', '−5.6%','Misleading — rebin', AMBER, LAMBER),
    ('F4  Debt-to-Equity Ratio',             '10.0%','1.48','10.5%','+ 0.5%','Broadly correct', TEAL, OFFWHT),
    ('F5  Debt Service Coverage Ratio',       '10.0%','1.66','11.8%','+ 1.8%','Broadly correct', TEAL, OFFWHT),
]
for i, (feat, ew, iv, iw, gap, verdict, vcol, rbg) in enumerate(rows):
    rect(sl, 0.45, 2.24+i*0.42, sum(ws)+0.3, 0.42, fill=rbg, border=LGRAY, border_w=0.3)
    vals = [feat, ew, iv, iw, gap, verdict]
    xx = 0.55
    for j, (v, w) in enumerate(zip(vals, ws)):
        c = vcol if j==5 else (DKGRAY if j!=4 else (RED if '−' in v else GREEN))
        bold = j in [0, 5]
        tb(sl, v, xx, 2.26+i*0.42, w, 0.38, sz=10, color=c, bold=bold)
        xx += w

divider(sl, 0.45, 5.62)
tb(sl, 'Recommended revision:  C (Credit History) → 55%  with  C1: 20%  |  C2: 10%  |  C3: 70%     F (Financial) → 45%  with  F1: 30%  |  F2: 15%  |  F3: Rebin  |  F4: 25%  |  F5: 30%',
   0.45, 5.7, 12.4, 0.5, sz=10, bold=True, color=NAVY)
tb(sl, 'Note: F3 Sale Growth is not removed outright — data shows inverted pattern suggesting rebinning and weight reduction before any exclusion decision.',
   0.45, 6.2, 12.4, 0.3, sz=9, italic=True, color=MGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 06 — Q1: SUPPORTING EVIDENCE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'Q1: Supporting Evidence — Default Rates by Feature',
           'Univariate default rates confirm IV rankings — C3 dominates; F3 shows no directional signal')

# KPI strip
kpis = [
    ('C3: Overdue >90d',     '100%',  RED),
    ('C1: Restr. <6 months', '62.5%', RED),
    ('F5: DSCR (0, 0.25)',   '18.1%', AMBER),
    ('F3: Growth >30%',      '11.4%', MGRAY),
    ('F3: Growth <=0%',      ' 9.1%', MGRAY),
    ('Overall Portfolio',    ' 9.0%', TEAL),
]
for i, (lbl, val, col) in enumerate(kpis):
    kpi(sl, 0.45+i*2.15, 1.18, 2.0, 1.05, val, lbl, val_color=col, val_sz=26)

divider(sl, 0.45, 2.36)

# C3 detail
bullet_card(sl, 0.45, 2.5, 5.9, 2.15, 'C3 — Severely Under-weighted (IV = 5.43)', [
    '▸  Strongest predictor by far — IV 4.6× higher than next feature (F1)',
    '▸  "Overdue >90d >1yr": 100% default rate across all 7 customers',
    '▸  "No delinquency": 5.6% default — 18× safer than worst bin',
    '▸  Chi-square: p < 0.001  |  Clear monotone risk ordering',
], sz=11, gap=0.44)

bullet_card(sl, 6.6, 2.5, 6.3, 2.15, 'F3 — Misleading Signal (IV = 0.62)', [
    '▸  Chi-square p = 0.758 — no statistically significant link to default',
    '▸  Spearman ρ = +0.696 — directionally WRONG (more growth = more default)',
    '▸  High growth firms (>30%) default at 11.4% vs 9.1% for stagnant firms',
    '▸  Possible explanation: aggressive growth may signal financial distress',
], sz=11, gap=0.44)

# Additional evidence row
bullet_card(sl, 0.45, 4.82, 5.9, 2.0, 'F4 Debt-to-Equity — Wrong Direction', [
    '▸  (0, 0.5) receives max score (+100 pts) but has 20% default rate',
    '▸  (0.5, 1) has 0% default rate but receives only +80 pts',
    '▸  Very low D/E may indicate under-leveraged, cash-poor SMEs',
    '▸  Bin redesign required before weight adjustment',
], sz=11, gap=0.41)

bullet_card(sl, 6.6, 4.82, 6.3, 2.0, 'Additional Validation — LR Coefficients', [
    '▸  Logistic regression on WoE features confirms C3 is top predictor',
    '▸  C3 coefficient: 0.890 (24.4% implied weight) vs expert 25.0% ✓',
    '▸  F4, F5 coefficients suggest higher weight than current 10% each',
    '▸  Caution: LR coefficients are log-odds based, not linear importance',
], sz=11, gap=0.41)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 07 — Q2: SCORECARD ENHANCEMENTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'Q2: Scorecard Enhancement Suggestions',
           'Six issues identified through monotonicity analysis, WoE calibration, and chi-square significance testing')

issues = [
    (RED,   'F3 Sale Growth\nInverted Scoring Direction',
     'Default rates are flat (6.8%–11.4%) with no monotone pattern.\nHigh growth firms default more than modest-growth firms.',
     'Rebin: consider U-shaped bins or significantly reduce weight.\nDo not assign +100 pts to highest-growth bin.'),
    (RED,   'F4 Debt-to-Equity\nPeak Bin Scoring Error',
     '(0, 0.5) assigned +100 pts despite 20% default rate.\n(0.5, 1) has 0% defaults but receives only +80 pts.',
     'Shift maximum score to (0.5, 1) bin.\nRedesign bin boundaries around empirical sweet spot.'),
    (AMBER, 'F2 Net Profit Margin\nNon-Monotone at Extremes',
     '>15% margin defaults at 12.0% — same rate as loss-making firms.\nModerate profitability (7.5–15%) is actually safest (0% default).',
     'Cap maximum score at (7.5, 15) bin.\nReduce score for >15% to reflect actual default rate.'),
    (AMBER, 'C3 Credit Bureau\nAmbiguous "Not Performed" Rule',
     'Two entries: "Not performed (exempted)" = +100 pts, 0% default.\n"Not performed" = −100 pts, 16.7% default.\n200-point gap for visually identical labels.',
     'Document a clear, auditable business rule to distinguish the two.\nConsider merging if the distinction cannot be operationalised.'),
    (AMBER, 'F5 DSCR\nMissing Data Scored as Neutral',
     '122 customers with missing DSCR default at 14.8% —\nnearly as risky as worst bin — yet receive 0 points (neutral).',
     'Assign a penalty score comparable to the (0, 0.25) bin.\nMissing DSCR is a meaningful risk signal, not neutral data.'),
    (TEAL,  'All Features\nAdopt PDO Score Calibration',
     'Current scores use arbitrary round numbers (−100, 0, +75, +100).\nPDO framework: Points = −(Factor × β × WoE)\nExample: C3 "Overdue >90d" should score −608 vs current −100.',
     'Base=600, PDO=20, Odds=50:1  →  Factor=28.85, Offset=487.1\nProvides auditable, mathematically grounded scores for regulators.'),
]
for i, (col, title, problem, fix) in enumerate(issues):
    c = i%2; r = i//2
    l = 0.45 + c*6.45
    t = 1.2  + r*2.05
    rect(sl, l, t, 6.2, 1.92, fill=WHITE, border=LGRAY, border_w=0.5)
    rect(sl, l, t, 0.08, 1.92, fill=col)
    rect(sl, l+0.08, t, 6.12, 0.36, fill=OFFWHT)
    tb(sl, title.replace('\n',' — '), l+0.2, t+0.05, 5.9, 0.28,
       sz=11, bold=True, color=DKGRAY)
    tb(sl, problem, l+0.2, t+0.42, 5.9, 0.72, sz=9, color=DKGRAY)
    rect(sl, l+0.08, t+1.17, 6.12, 0.03, fill=LGRAY)
    tb(sl, '→  ' + fix, l+0.2, t+1.25, 5.9, 0.6,
       sz=9, bold=True, color=col if col!=AMBER else GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 08 — Q3: EDA METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'Q3: EDA Methodology for Default Prediction',
           'Generalised framework applicable to any credit risk dataset with financial, credit history, and macroeconomic variables')

# Process flow arrows
flow = ['Data Quality\n& Profiling', 'Univariate\nAnalysis', 'Bivariate\nAnalysis', 'WoE / IV\nAnalysis']
for i, label in enumerate(flow):
    l = 0.45 + i*3.25
    rect(sl, l, 1.18, 2.85, 0.62, NAVY if i==3 else TEAL)
    tb(sl, f'STEP {i+1}', l+0.12, 1.2, 2.6, 0.22, sz=8, bold=True, color=GOLD)
    tb(sl, label, l+0.12, 1.38, 2.6, 0.38, sz=11, bold=True, color=WHITE)
    if i < 3:
        tb(sl, '→', l+2.85+0.02, 1.32, 0.38, 0.36, sz=18, bold=True, color=MGRAY)

step_content = [
    [
        '▸  Identify variable types: categorical, numerical, time-series macro',
        '▸  Assess class imbalance — rare events require balanced modelling',
        '▸  Distinguish true nulls from sentinel values (encode separately)',
        '▸  Check outliers using IQR — more robust than SD for skewed data',
        '▸  Skewness: guides bin placement; WoE encoding neutralises raw skew',
    ],[
        '▸  Categorical: default rate + concentration per bin',
        '▸  Numerical: boxplots / histograms split by default flag',
        '▸  Monotonicity: default rate must move consistently across bins',
        '▸  Macro: time-series plot vs portfolio default rate (lagged 1–4Q)',
        '▸  Concentration check: bins with <5% population → unstable WoE',
    ],[
        '▸  Numerical vs default: point-biserial correlation + KS statistic',
        '▸  Categorical vs default: chi-square test (p < 0.05 threshold)',
        '▸  Multicollinearity: Pearson matrix + Cramér\'s V for categoricals',
        '▸  Macro: Granger causality test — does past macro predict default?',
        '▸  Vintage analysis: default rate by origination year → macro effects',
    ],[
        '▸  WoE_i = ln(% Bad_i / % Good_i)  →  already in log-odds space',
        '▸  IV = Σ (% Bad − % Good) × WoE  →  ranks feature importance',
        '▸  IV > 0.30: strong predictor | IV < 0.02: exclude',
        '▸  Non-monotone WoE → merge bins before scorecard inclusion',
        '▸  Build WoE maps on training data only — avoid data leakage',
    ],
]

for i, items in enumerate(step_content):
    l = 0.45 + (i%2)*6.45
    t = 1.98  + (i//2)*2.6
    card(sl, l, t, 6.2, 2.42, title=flow[i].replace('\n',' '),
         title_bg=NAVY if i==3 else TEAL)
    for j, item in enumerate(items):
        tb(sl, item, l+0.15, t+0.42+j*0.4, 5.9, 0.38, sz=10, color=DKGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 09 — Q4: PD MODEL DEVELOPMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'Q4: Logistic Regression for PD Estimation',
           'Theory, development steps, and underlying assumptions')

# Theory strip
rect(sl, 0.45, 1.18, 12.45, 0.62, OFFWHT, border=LGRAY, border_w=0.5)
tb(sl, 'Why Logistic Regression?', 0.6, 1.2, 3.5, 0.26, sz=10, bold=True, color=NAVY)
tb(sl, 'Linear regression predicts unbounded values — PD must lie in [0,1]. Logistic regression applies the sigmoid to map any input to a valid probability:',
   0.6, 1.42, 7.5, 0.34, sz=9, italic=True, color=DKGRAY)
tb(sl, 'logit(PD) = ln(PD/1−PD) = β₀ + β₁·WoE₁ + ... + βₙ·WoEₙ     →     PD = 1 / (1 + e^−z)',
   8.3, 1.28, 4.5, 0.46, sz=10, bold=True, color=TEAL)

# Six steps
steps = [
    ('1. Data Preparation',
     '▸  Handle sentinel values as separate bins\n▸  Fill nulls → "Not available" category\n▸  Define target: Default Flag (0/1)'),
    ('2. WoE Encoding',
     '▸  Replace bins with log-odds values\n▸  Linearises LR assumption by construction\n▸  Comparable scale — no standardisation needed'),
    ('3. Train / Test Split',
     '▸  Stratified 80/20 split — preserves 9% default rate\n▸  5-fold stratified CV preferred given only 50 defaults\n▸  Build WoE maps on train set only (avoid leakage)'),
    ('4. Model Training',
     '▸  class_weight="balanced" — corrects 9:91 imbalance\n▸  L2 regularisation — tune C via grid search\n▸  Coefficients should all be positive (WoE encodes direction)'),
    ('5. PDO Calibration',
     '▸  Points_i = −(Factor × β_i × WoE_i)\n▸  Base=600, PDO=20, Odds=50:1\n▸  Factor=28.85, Offset=487.1'),
    ('6. Score → PD',
     '▸  PD = 1 / (1 + e^((score−offset)/factor))\n▸  Every +20 pts → odds double → PD halves\n▸  Interpretable probability for credit decisions'),
]
for i, (title, body) in enumerate(steps):
    c = i%3; r = i//2
    l = 0.45 + c*4.3
    t = 1.98 + r*1.88
    rect(sl, l, t, 4.0, 1.74, fill=WHITE, border=LGRAY, border_w=0.5)
    rect(sl, l, t, 4.0, 0.34, fill=TEAL if r==0 else NAVY)
    tb(sl, title, l+0.12, t+0.05, 3.8, 0.28, sz=11, bold=True, color=WHITE)
    tb(sl, body,  l+0.12, t+0.4,  3.8, 1.28, sz=9.5, color=DKGRAY)

# Assumptions footer
divider(sl, 0.45, 5.76)
tb(sl, 'Key assumptions: Binary Bernoulli outcome  |  Linearity in log-odds (satisfied by WoE)  |  Independence of observations  |  No perfect multicollinearity  |  Sufficient EPV (≥10)',
   0.45, 5.82, 12.4, 0.3, sz=9, color=MGRAY)
tb(sl, 'Critical limitation: EPV = 50 defaults / 8 features = 6.25 < 10 → coefficients are unstable; hybrid expert-statistical approach mitigates this.',
   0.45, 6.12, 12.4, 0.3, sz=9, italic=True, color=RED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — MODEL RESULTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'Model Results',
           'Logistic regression on WoE-encoded features — strong discriminatory power')

kpi(sl, 0.45, 1.18, 2.9, 1.05, '0.87',  'AUC-ROC',            GREEN,  30)
kpi(sl, 3.55, 1.18, 2.9, 1.05, '0.74',  'Gini Coefficient',   GREEN,  30)
kpi(sl, 6.65, 1.18, 2.9, 1.05, '0.82',  '5-Fold CV AUC',      TEAL,   30)
kpi(sl, 9.75, 1.18, 3.1, 1.05, '+441',  'Good vs Bad Score Gap (pts)', NAVY, 30)

divider(sl, 0.45, 2.36)

# LR coefficient table (left)
tb(sl, 'LR-Implied Feature Importance', 0.45, 2.48, 6.0, 0.3,
   sz=11, bold=True, color=NAVY)
tb(sl, 'Note: coefficients are log-odds based — IV remains primary basis for weight justification',
   0.45, 2.75, 6.0, 0.22, sz=8, italic=True, color=MGRAY)

coef_data = [
    ('C3  Credit Bureau',          '0.890', '24.4%', GREEN,  LGREEN),
    ('F4  Debt-to-Equity',         '0.754', '20.7%', GREEN,  LGREEN),
    ('F5  DSCR',                   '0.583', '16.0%', GREEN,  LGREEN),
    ('C2  Debt Rescheduling',      '0.312', ' 8.6%', DKGRAY, OFFWHT),
    ('C1  Debt Restructuring',     '0.308', ' 8.4%', DKGRAY, OFFWHT),
    ('F3  Sale Growth',            '0.297', ' 8.1%', AMBER,  LAMBER),
    ('F1  Net Fixed Asset',        '0.285', ' 7.8%', DKGRAY, OFFWHT),
    ('F2  Net Profit Margin',      '0.220', ' 6.0%', DKGRAY, OFFWHT),
]
rect(sl, 0.45, 3.0, 6.1, 0.32, NAVY)
for j, (h, w) in enumerate(zip(['Feature','Coefficient','Implied Weight'],
                                 [3.3, 1.4, 1.4])):
    xx = 0.55 + sum([3.3,1.4,1.4][:j])
    tb(sl, h, xx, 3.02, w, 0.28, sz=9, bold=True, color=WHITE)

for i, (feat, coef, wt, col, bg) in enumerate(coef_data):
    rect(sl, 0.45, 3.32+i*0.38, 6.1, 0.38, fill=bg, border=LGRAY, border_w=0.3)
    tb(sl, feat, 0.55, 3.34+i*0.38, 3.25, 0.32, sz=9.5, color=DKGRAY)
    tb(sl, coef, 3.85, 3.34+i*0.38, 1.35, 0.32, sz=9.5, color=DKGRAY)
    tb(sl, wt,   5.25, 3.34+i*0.38, 1.25, 0.32, sz=9.5, bold=True, color=col)

# Score distribution (right)
bullet_card(sl, 6.75, 2.48, 6.1, 2.18, 'Score Distribution by Default Status', [
    '▸  Good customers (Default=0):  Mean = 785 pts  |  Median = 551 pts',
    '▸  Bad customers  (Default=1):  Mean = 344 pts  |  Median = 454 pts',
    '▸  Score gap of 441 pts — strong separation between groups',
    '▸  Score range: −335 (riskiest) to +2,321 (safest)',
], sz=10.5, gap=0.44)

bullet_card(sl, 6.75, 4.82, 6.1, 2.0, 'Score → PD Interpretation (PDO Scale)', [
    '▸  Score 645  →  PD ≈ 1.0%   (very low risk)',
    '▸  Score 600  →  PD ≈ 2.0%   (base: 50:1 good:bad odds)',
    '▸  Score 572  →  PD ≈ 5.0%   (moderate risk)',
    '▸  Score 511  →  PD ≈ 20.0%  (high risk)',
], sz=10.5, gap=0.41)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Q5: MODEL EVALUATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl)
header_bar(sl, 'Q5: Model Performance Evaluation',
           'Three pillars: Discriminatory Power  |  Goodness-of-Fit  |  Forecasting Accuracy & Stability')

pillars = [
    ('Discriminatory Power', TEAL, [
        ('AUC-ROC', 'Area under ROC curve. Target > 0.75.\nModel AUC = 0.87 ✓'),
        ('Gini Coefficient', '2 × AUC − 1 = 0.74.\nMeasures rank-ordering ability.'),
        ('KS Statistic', 'Max separation between good and\nbad cumulative distributions.\nTarget KS > 0.30.'),
        ('Lorenz Curve', 'Visual representation of how well\nthe model concentrates defaults\nin the lower score deciles.'),
    ]),
    ('Goodness of Fit', NAVY, [
        ('Log-Likelihood', 'How well the model fits observed\ndefaults. Higher is better.'),
        ('AIC / BIC', 'Penalises complexity — use to\ncompare models with different\nlag structures or feature sets.'),
        ('Hosmer-Lemeshow', 'Tests if predicted PD equals\nactual default rate across\ndecile groups. p > 0.05 = pass.'),
        ('Calibration Curve', 'Plots predicted PD vs observed\ndefault rate — checks if PD\nestimates are well-calibrated.'),
    ]),
    ('Forecasting Accuracy\n& Stability', AMBER, [
        ('Confusion Matrix', 'TP, FP, FN, TN — basis for\nsensitivity and precision metrics.\nTarget sensitivity ≥ 80%.'),
        ('PSI (Stability)', 'Population Stability Index.\nPSI < 0.10 = stable distribution.\nPSI > 0.20 = significant drift.'),
        ('Backtesting', 'Compare predicted PD to realised\ndefault rates over time.\nRequired under Basel II / IFRS 9.'),
        ('Score Distribution', 'Monitor score distributions over\ntime — flag shifts in customer\nrisk profile or data drift.'),
    ]),
]
for i, (title, col, metrics) in enumerate(pillars):
    l = 0.45 + i*4.3
    rect(sl, l, 1.18, 4.0, 5.85, fill=WHITE, border=LGRAY, border_w=0.5)
    rect(sl, l, 1.18, 4.0, 0.42, fill=col)
    tb(sl, title, l+0.14, 1.21, 3.8, 0.36, sz=12, bold=True, color=WHITE)
    for j, (name, desc) in enumerate(metrics):
        t = 1.72 + j*1.3
        rect(sl, l+0.14, t, 3.72, 1.12, fill=OFFWHT, border=LGRAY, border_w=0.3)
        tb(sl, name, l+0.26, t+0.06, 3.5, 0.26, sz=10, bold=True, color=DKGRAY)
        tb(sl, desc,  l+0.26, t+0.32, 3.5, 0.74, sz=9,  color=MGRAY)

divider(sl, 0.45, 7.08)
tb(sl, 'Validation note: With only 50 defaults, use 5-fold stratified CV + bootstrap confidence intervals for AUC. A single 80/20 split is insufficient for reliable performance estimation.',
   0.45, 7.15, 12.4, 0.3, sz=9, italic=True, color=MGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION & RECOMMENDATIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_bg(sl, NAVY)
rect(sl, 0, 0, 0.5, 7.5, TEAL)
rect(sl, 0, 6.45, 13.33, 0.06, GOLD)
tb(sl, 'Conclusion & Recommendations', 0.7, 0.12, 12, 0.65,
   sz=26, bold=True, color=WHITE)
tb(sl, 'Actionable findings to strengthen the SME credit scorecard',
   0.7, 0.72, 12, 0.35, sz=12, italic=True, color=GOLD)

recs = [
    (TEAL,  '01',
     'Reweight C3 Credit Bureau',
     'Increase weight from 25% to ~38% based on IV=5.43.\nC3 is the single most predictive feature — 4× the IV of the next best.'),
    (TEAL,  '02',
     'Reweight & Rebin F3 Sale Growth',
     'Reduce weight significantly and rebin.\nCurrent scoring direction is inverted — aggressive growth correlates with higher default.'),
    (AMBER, '03',
     'Fix F4 Debt-to-Equity Scoring',
     'Redesign bins so (0.5, 1) receives maximum score.\nCurrent peak bin (0, 0.5) has 20% default rate — the worst outcome.'),
    (AMBER, '04',
     'Adopt PDO Score Calibration',
     'Replace arbitrary expert scores with PDO-calibrated points.\nProvides auditable, mathematically grounded output for regulatory review.'),
    (RED,   '05',
     'Address EPV Constraint',
     '50 defaults / 8 features = 6.25 EPV (below threshold of 10).\nReduce features, increase regularisation (C=0.1), use 5-fold CV.'),
    (GREEN, '06',
     'Extend to Macro Variable Integration',
     'Dataset lacks time dimension — recommend adding GDP growth (lag 4Q)\nand interest rate as overlay for IFRS 9 forward-looking stress testing.'),
]
for i, (col, num, title, body) in enumerate(recs):
    c = i%3; r = i//2
    l = 0.65 + c*4.22
    t = 1.22 + r*2.5
    rect(sl, l, t, 3.9, 2.22, fill=RGBColor(0x10, 0x28, 0x48),
         border=col, border_w=1.5)
    rect(sl, l, t, 0.55, 2.22, fill=col)
    tb(sl, num, l+0.06, t+0.8, 0.44, 0.6, sz=18, bold=True, color=WHITE,
       align=PP_ALIGN.CENTER)
    tb(sl, title, l+0.65, t+0.1, 3.17, 0.52, sz=11, bold=True, color=WHITE)
    tb(sl, body,  l+0.65, t+0.64, 3.17, 1.5,  sz=9.5, color=LGRAY)

tb(sl, 'The hybrid expert-statistical approach is the appropriate framework given the limited default sample. Data analysis provides the objective evidence; expert judgement provides the stable anchor.',
   0.65, 6.55, 12, 0.52, sz=9.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────────
path = '/Users/leekim/prj/deloitte/Credit_Scorecard_Presentation_v2.pptx'
prs.save(path)
print(f'Saved: {path}')
print(f'Slides: {len(prs.slides)}')
