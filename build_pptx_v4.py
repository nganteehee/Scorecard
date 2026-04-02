"""
build_pptx_v4.py  –  Polished professional presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)
NAVY2  = RGBColor(0x0A, 0x1F, 0x40)   # darker navy for depth
TEAL   = RGBColor(0x00, 0x7A, 0x8A)
LTEAL  = RGBColor(0xE4, 0xF4, 0xF6)   # light teal bg
GOLD   = RGBColor(0xE8, 0xA0, 0x00)
LGOLD  = RGBColor(0xFD, 0xF3, 0xDC)   # light gold bg
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF4, 0xF8)
MGRAY  = RGBColor(0x88, 0x99, 0xAA)
DGRAY  = RGBColor(0x44, 0x55, 0x66)
RED    = RGBColor(0xC0, 0x39, 0x2B)
LRED   = RGBColor(0xFD, 0xED, 0xEB)
GREEN  = RGBColor(0x1A, 0x8A, 0x50)
LGREEN = RGBColor(0xE8, 0xF8, 0xEE)
AMBER  = RGBColor(0xD4, 0x7E, 0x00)

IMG  = '/Users/leekim/prj/deloitte/'
W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def R(slide, l, t, w, h, fill, line_color=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def T(slide, text, l, t, w, h, size=11, bold=False, color=NAVY,
      align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def PIC(slide, path, l, t, w, h=None):
    full = IMG + path if not path.startswith('/') else path
    if not os.path.exists(full): return
    if h:
        slide.shapes.add_picture(full, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(full, Inches(l), Inches(t), width=Inches(w))

def header(slide, title, subtitle=None):
    """Navy header with gold accent line and optional subtitle."""
    R(slide, 0, 0, 13.33, 1.15, NAVY)
    R(slide, 0, 1.15, 13.33, 0.055, GOLD)
    T(slide, title,    0.45, 0.10, 12.4, 0.60, size=21, bold=True, color=WHITE)
    if subtitle:
        T(slide, subtitle, 0.45, 0.72, 12.4, 0.38, size=10, color=GOLD, italic=True)

def footer(slide, page_num, total=14):
    """Thin footer bar with slide number and branding."""
    R(slide, 0, 7.28, 13.33, 0.22, NAVY)
    T(slide, 'SME Credit Scorecard Development  |  Deloitte Credit Risk',
      0.4, 7.29, 10.0, 0.20, size=7.5, color=MGRAY)
    T(slide, f'{page_num} / {total}', 12.0, 7.29, 1.25, 0.20,
      size=7.5, color=MGRAY, align=PP_ALIGN.RIGHT)

def bg(slide):
    R(slide, 0, 0, 13.33, 7.5, LGRAY)

def divider(slide, t, l=0.45, w=12.43, color=TEAL, thick=0.03):
    R(slide, l, t, w, thick, color)

def kpi(slide, l, t, w, h, value, label, sub='', vcolor=TEAL, tag_color=TEAL):
    """KPI card: colored top strip, big value, label, optional subtext."""
    R(slide, l, t, w, h, WHITE, line_color=LGRAY)
    R(slide, l, t, w, 0.07, tag_color)
    T(slide, value, l+0.1, t+0.12, w-0.2, 0.52,
      size=24, bold=True, color=vcolor, align=PP_ALIGN.CENTER)
    T(slide, label, l+0.1, t+0.66, w-0.2, 0.28,
      size=8.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    if sub:
        T(slide, sub, l+0.1, t+0.93, w-0.2, 0.22,
          size=7.5, color=MGRAY, align=PP_ALIGN.CENTER)

def card_left(slide, l, t, w, h, title, bullets, accent=TEAL, bg_color=WHITE):
    """Card with left accent bar."""
    R(slide, l, t, w, h, bg_color, line_color=LGRAY)
    R(slide, l, t, 0.055, h, accent)
    T(slide, title, l+0.18, t+0.1, w-0.28, 0.3,
      size=9.5, bold=True, color=NAVY)
    divider(slide, t+0.42, l=l+0.18, w=w-0.3, color=LGRAY, thick=0.02)
    y = t + 0.5
    for b in bullets:
        T(slide, f'• {b}', l+0.18, y, w-0.28, 0.28, size=8, color=DGRAY)
        y += 0.31

def card_top(slide, l, t, w, h, title, bullets, accent=TEAL):
    """Card with top accent bar and title."""
    R(slide, l, t, w, h, WHITE, line_color=LGRAY)
    R(slide, l, t, w, 0.42, accent)
    T(slide, title, l+0.15, t+0.08, w-0.3, 0.3,
      size=9.5, bold=True, color=WHITE)
    y = t + 0.52
    for b in bullets:
        T(slide, f'• {b}', l+0.15, y, w-0.3, 0.3, size=8, color=DGRAY)
        y += 0.31

def tag(slide, text, l, t, w, color=TEAL, text_color=WHITE):
    """Small pill-shaped label."""
    R(slide, l, t, w, 0.27, color)
    T(slide, text, l+0.05, t+0.04, w-0.1, 0.20,
      size=7.5, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def section_row(slide, l, t, w, h, num, title, detail, num_color=TEAL):
    """Numbered section row for tables."""
    R(slide, l, t, 0.55, h, num_color)
    T(slide, num, l+0.05, t+(h-0.3)/2, 0.45, 0.3,
      size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(slide, title, l+0.65, t+0.06, w*0.35, h-0.12,
      size=8.5, bold=True, color=NAVY)
    T(slide, detail, l+0.65+w*0.35+0.1, t+0.06, w*0.58, h-0.12,
      size=8, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
R(sl, 0, 0, 13.33, 7.5, NAVY2)
R(sl, 0, 0, 0.42, 7.5, TEAL)          # left teal bar
R(sl, 0.42, 0, 12.91, 0.08, TEAL)     # top teal strip
R(sl, 0.42, 7.42, 12.91, 0.08, GOLD)  # bottom gold strip

# decorative circles
R(sl, 10.5, -1.0, 4.5, 4.5, RGBColor(0x10, 0x35, 0x68))
R(sl, 11.2, 5.0,  3.5, 3.5, RGBColor(0x10, 0x35, 0x68))

R(sl, 0.85, 2.9, 11.6, 0.06, GOLD)   # gold divider

T(sl, 'SME Credit Scorecard Development',
  1.0, 1.1, 11.0, 1.3, size=38, bold=True, color=WHITE)
T(sl, 'Credit Risk Model: Weight Assessment, EDA & PD Estimation',
  1.0, 2.55, 11.0, 0.55, size=15, color=GOLD)
T(sl, 'Candidate Technical Presentation  |  Deloitte Credit Risk',
  1.0, 3.18, 11.0, 0.4, size=11, color=MGRAY)

# stats row
for i, (val, lbl) in enumerate([('557','Borrowers'),('8','Features'),
                                  ('9.0%','Default Rate'),('0.90','AUC-ROC')]):
    x = 1.0 + i * 2.8
    R(sl, x, 4.1, 2.4, 0.85, RGBColor(0x10, 0x38, 0x70))
    R(sl, x, 4.1, 2.4, 0.05, GOLD)
    T(sl, val, x+0.1, 4.18, 2.2, 0.42, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, lbl, x+0.1, 4.60, 2.2, 0.28, size=8,  color=GOLD, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'Executive Summary',
               'Key findings across weight assessment, model performance and recommendations')
footer(sl, 2)

# KPI row
kpis = [
    ('0.90', 'AUC-ROC',          'Excellent discrimination', TEAL,  TEAL),
    ('0.80', 'Gini Coefficient',  'Strong predictive power',  TEAL,  TEAL),
    ('9.0%', 'Portfolio Default', '50 / 557 borrowers',       NAVY,  NAVY),
    ('C3',   'Top Predictor',     'CB Delinquency IV = 5.43', TEAL,  TEAL),
    ('D/E',  'Underweighted',     'LR 22% vs Expert 10%',     AMBER, AMBER),
]
x = 0.45
for val, lbl, sub, vc, tc in kpis:
    kpi(sl, x, 1.28, 2.44, 1.25, val, lbl, sub, vc, tc)
    x += 2.55

divider(sl, 2.72)

# Three finding cards
findings = [
    ('Q1 — Weight Assessment', TEAL,
     ['C3 CB Delinquency IV=5.43 — dominant predictor',
      'D/E Ratio underweighted: expert 10% vs LR-implied 22%',
      'DSCR underweighted: expert 10% vs LR-implied 16%',
      'F3 Sale Growth IV=0.04 — negligible signal, remove']),
    ('Q2 — Scorecard Enhancements', NAVY,
     ['Apply Laplace smoothing to avoid WoE explosion (±20)',
      'Merge zero-default bins to improve monotonicity',
      'Recalibrate PD — class_weight inflates PD 4×',
      'Add macroeconomic overlay (GDP, interest rates)']),
    ('Q4/Q5 — Model & Validation', TEAL,
     ['LR with WoE encoding, AUC = 0.90, Gini = 0.80',
      'HL test fails: extreme WoE distorts calibration',
      'Binomial test: predicted 35% vs actual 9%',
      'Score mapped to 300–850; clear band separation']),
]
x = 0.45
for title, accent, bullets in findings:
    card_left(sl, x, 2.82, 4.1, 4.25, title, bullets, accent=accent)
    x += 4.27

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Background & Data Overview
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'I. Background & Data Overview',
               'SME portfolio scorecard — expert-weighted system under review')
footer(sl, 3)

# Left panel — context
R(sl, 0.45, 1.28, 5.9, 5.75, WHITE, line_color=LGRAY)
R(sl, 0.45, 1.28, 5.9, 0.42, NAVY)
T(sl, 'Dataset Context', 0.65, 1.33, 5.5, 0.32, size=10, bold=True, color=WHITE)
context = [
    ('Dataset',      '557 SME borrowers · 8 features · 9.0% default rate'),
    ('Target',       'Default Flag  (1 = default, 0 = performing)'),
    ('Structure',    'C component 50% + F component 50%'),
    ('C Component',  'C1 Restructuring · C2 Rescheduling · C3 CB Delinquency'),
    ('F Component',  'F1–F5 financial ratios, each weighted 10%'),
    ('Missing',      'C2: 282 NaN (51%) · Net Fixed Asset: 26 rows (5%)'),
    ('Sentinels',    '99999999 = zero denominator · 99999998 = negative equity'),
    ('Objective',    'Assess expert weights via IV; build LR-based PD model'),
]
y = 1.82
for k, v in context:
    R(sl, 0.47, y-0.02, 5.86, 0.36, LGRAY if context.index((k,v)) % 2 == 0 else WHITE)
    T(sl, k, 0.62, y+0.02, 1.55, 0.28, size=8.5, bold=True, color=TEAL)
    T(sl, v, 2.22, y+0.02, 3.98, 0.28, size=8.5, color=DGRAY)
    y += 0.38

# Right panel — feature table
R(sl, 6.6, 1.28, 6.3, 5.75, WHITE, line_color=LGRAY)
R(sl, 6.6, 1.28, 6.3, 0.42, NAVY)
T(sl, 'Feature Summary', 6.8, 1.33, 5.9, 0.32, size=10, bold=True, color=WHITE)

# table headers
hdrs = ['Code', 'Feature', 'Type', 'IV', 'Expert Wt', 'Signal']
xs   = [6.65,  7.18,  9.42,  10.48, 11.18, 11.88]
ws   = [0.48,  2.19,  1.0,   0.65,  0.65,  0.97]
y = 1.80
for i, h in enumerate(hdrs):
    T(sl, h, xs[i], y+0.48, ws[i], 0.28, size=8, bold=True, color=NAVY)
divider(sl, 2.18, l=6.65, w=6.2, color=TEAL, thick=0.025)

rows = [
    ('C1','Debt Restructuring Hist.','Cat','1.17','12.5%','Strong', TEAL),
    ('C2','Debt Rescheduling Hist.', 'Cat','0.57','12.5%','Moderate',TEAL),
    ('C3','Worst CB Delinquency 2yr','Cat','5.43','25.0%','Strong', TEAL),
    ('F1','Net Fixed Asset',         'Num','2.60','10.0%','Strong', NAVY),
    ('F2','Net Profit Margin (%)',   'Num','1.28','10.0%','Strong', NAVY),
    ('F3','Sale Growth (%)',         'Num','0.04','10.0%','Weak',   RED),
    ('F4','Debt-to-Equity Ratio',    'Num','1.46','10.0%','Strong', NAVY),
    ('F5','Debt Service Coverage',   'Num','1.58','10.0%','Strong', NAVY),
]
y = 2.28
for r in rows:
    bg2 = LGRAY if rows.index(r) % 2 == 0 else WHITE
    R(sl, 6.62, y-0.02, 6.25, 0.34, bg2)
    code, feat, typ, iv, ew, sig, cc = r
    signal_col = RED if sig == 'Weak' else GREEN
    for i, v in enumerate([code, feat, typ, iv, ew]):
        T(sl, v, xs[i], y+0.01, ws[i], 0.28, size=8,
          color=RED if v == '0.04' else DGRAY,
          bold=(i == 0))
    tag(sl, sig, xs[5], y, ws[5]-0.05,
        color=LRED if sig == 'Weak' else LGREEN,
        text_color=RED if sig == 'Weak' else GREEN)
    y += 0.36

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — EDA: Univariate Analysis
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'II. EDA: Univariate Analysis',
               'Distribution and default rates across categorical and numerical features')
footer(sl, 4)

PIC(sl, 'eda_univariate_boxplot.png', 0.45, 1.28, 8.3, 4.6)

R(sl, 8.95, 1.28, 4.0, 4.6, WHITE, line_color=LGRAY)
R(sl, 8.95, 1.28, 4.0, 0.42, NAVY)
T(sl, 'Key Observations', 9.12, 1.33, 3.7, 0.32, size=10, bold=True, color=WHITE)
obs = [
    ('C1 Restructuring',
     '6–12 month restructured debt: 44% default vs 8% baseline — clear risk cliff'),
    ('C3 CB Delinquency',
     '>90-day overdue (>1yr ago): 100% default rate — strongest single predictor'),
    ('Net Profit Margin',
     'Heavy left-skew; median = −8%; majority of SMEs are unprofitable'),
    ('Sale Growth (%)',
     'Extreme right skew (skewness = 10.95); near-zero signal in default separation'),
    ('DSCR',
     'Defaulters cluster at DSCR ≤ 0; non-defaulters spread across positive range'),
]
y = 1.82
for title, detail in obs:
    T(sl, title, 9.12, y, 3.7, 0.22, size=8.5, bold=True, color=TEAL)
    T(sl, detail, 9.12, y+0.22, 3.7, 0.38, size=8, color=DGRAY)
    y += 0.7

divider(sl, 6.05)
R(sl, 0.45, 6.1, 12.43, 1.05, WHITE, line_color=LGRAY)
stats = [
    ('C2 Missing', '282 / 557 (51%) — treat NaN as "No rescheduling history" bin'),
    ('Sentinel Values', '99999999 in Profit Margin (52 rows) and Sale Growth (68 rows) → separate bin'),
    ('Skewness Alert', 'Net Fixed Asset: 6.94 · Sale Growth: 10.95 · DSCR: 14.45 → IQR capping required'),
]
x = 0.55
for label, detail in stats:
    T(sl, label, x, 6.18, 2.2, 0.25, size=8, bold=True, color=NAVY)
    T(sl, detail, x, 6.42, 3.8, 0.55, size=7.8, color=DGRAY)
    x += 4.18

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EDA: Correlation & Multicollinearity
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'II. EDA: Correlation & Multicollinearity',
               "Pearson correlation for numerical features; Cramér's V for categorical pairs")
footer(sl, 5)

PIC(sl, 'eda_correlation.png', 0.45, 1.28, 6.5, 4.35)
PIC(sl, 'eda_cramers_v.png',   7.15, 1.28, 3.55, 2.95)

R(sl, 7.15, 4.38, 5.73, 2.65, WHITE, line_color=LGRAY)
R(sl, 7.15, 4.38, 5.73, 0.42, NAVY)
T(sl, "Cramér's V — Categorical Multicollinearity", 7.32, 4.43, 5.4, 0.32,
  size=9.5, bold=True, color=WHITE)
cramers = [
    'C1 ↔ C2 = 0.40  |  C1 ↔ C3 = 0.37  |  C2 ↔ C3 = 0.39',
    'Moderate inter-correlation among all three credit history features',
    'Including all three may introduce partial redundancy in the model',
    'Recommendation: retain all — V < 0.60 (below critical threshold)',
]
y = 4.9
for c in cramers:
    T(sl, f'• {c}', 7.32, y, 5.42, 0.35, size=8.5, color=DGRAY)
    y += 0.42

divider(sl, 6.1)
R(sl, 0.45, 6.15, 6.45, 1.05, WHITE, line_color=LGRAY)
R(sl, 0.45, 6.15, 6.45, 0.42, TEAL)
T(sl, 'Numerical vs Default — KS Test', 0.62, 6.20, 6.1, 0.32,
  size=9.5, bold=True, color=WHITE)
ks_rows = [
    ('Net Fixed Asset', 'KS=0.30', 'p=0.0005', '✔ Significant'),
    ('DSCR',            'KS=0.37', 'p=0.0004', '✔ Significant'),
    ('Profit Margin',   'KS=0.18', 'p=0.113',  '✘ Not significant'),
    ('Sale Growth',     'KS=0.11', 'p=0.693',  '✘ Not significant'),
]
xk = [0.55, 2.3, 3.4, 4.5]
wk = [1.7,  1.05, 1.05, 1.75]
y = 6.67
for feat, ks, p, sig in ks_rows:
    sc = GREEN if '✔' in sig else RED
    for i, v in enumerate([feat, ks, p, sig]):
        T(sl, v, xk[i], y, wk[i], 0.27, size=8,
          color=sc if i == 3 else DGRAY, bold=(i == 3))
    y = 6.67  # single line — place side by side
xk2 = [0.55, 3.4, 5.8, 9.5]  # horizontal layout
for i2, (feat, ks, p, sig) in enumerate(ks_rows):
    sc = GREEN if '✔' in sig else RED
    T(sl, f'{feat}: {ks}  {p}  ', xk2[i2%4 if i2 < 4 else 0], 6.67, 3.1, 0.28,
      size=7.8, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — EDA: WoE Monotonicity
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'II. EDA: WoE Monotonicity Analysis',
               'Weight of Evidence validates risk ordering — bins must show a clear directional trend')
footer(sl, 6)

PIC(sl, 'eda_monotonicity.png', 0.45, 1.28, 12.43, 4.72)

divider(sl, 6.15)
notes = [
    (TEAL,  'C1 Restructuring',
     'Non-monotone: "Never restructured" WoE < "No history" WoE — these bins should be merged'),
    (AMBER, 'C3 CB Delinquency',
     'Broadly monotone; zero-default bins produce WoE = ±20 — Laplace smoothing required'),
    (RED,   'F4 D/E Ratio',
     '(0.5,1) bin has 0 defaults → WoE = −20.23 — extreme value distorts scoring'),
]
x = 0.45
for color, label, detail in notes:
    R(sl, x, 6.2, 4.1, 1.05, WHITE, line_color=LGRAY)
    R(sl, x, 6.2, 0.06, 1.05, color)
    T(sl, label,  x+0.18, 6.28, 3.75, 0.28, size=8.5, bold=True, color=NAVY)
    T(sl, detail, x+0.18, 6.56, 3.75, 0.55, size=8, color=DGRAY)
    x += 4.27

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Q1: Weight Assessment
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'III. Feature Weight Assessment',
               'Information Value (IV) as objective, data-driven measure vs expert-assigned weights')
footer(sl, 7)

PIC(sl, 'iv_comparison_chart.png', 0.45, 1.28, 7.9, 3.9)

# IV assessment table
R(sl, 8.55, 1.28, 4.35, 5.75, WHITE, line_color=LGRAY)
R(sl, 8.55, 1.28, 4.35, 0.42, NAVY)
T(sl, 'IV Assessment', 8.72, 1.33, 4.05, 0.32, size=10, bold=True, color=WHITE)
hdrs2 = ['Feature', 'Expert', 'IV', 'Verdict']
xs2   = [8.60,  9.95, 10.52, 11.1]
ws2   = [1.3,   0.52, 0.52,  1.25]
y = 1.82
for i, h in enumerate(hdrs2):
    T(sl, h, xs2[i], y, ws2[i], 0.28, size=8, bold=True, color=NAVY)
divider(sl, 2.13, l=8.6, w=4.25, color=TEAL)
iv_rows = [
    ('C3 CB Delinquency', '25%',  '5.43', '✅ Aligned',      GREEN, LGREEN),
    ('C1 Restructuring',  '12.5%','1.17', '✅ Reasonable',   GREEN, LGREEN),
    ('F1 Net Fixed Asset','10%',  '2.60', '⚠ Underweighted', AMBER, LGOLD),
    ('F4 D/E Ratio',      '10%',  '1.46', '⚠ Underweighted', AMBER, LGOLD),
    ('F5 DSCR',           '10%',  '1.58', '⚠ Underweighted', AMBER, LGOLD),
    ('F2 Profit Margin',  '10%',  '1.28', '⚠ Underweighted', AMBER, LGOLD),
    ('C2 Rescheduling',   '12.5%','0.57', '⚠ Overweighted',  AMBER, LGOLD),
    ('F3 Sale Growth',    '10%',  '0.04', '❌ Remove',        RED,   LRED),
]
y = 2.22
for feat, exp, iv, verdict, vc, bgc in iv_rows:
    R(sl, 8.57, y-0.02, 4.3, 0.33, bgc)
    T(sl, feat,    xs2[0], y, ws2[0], 0.28, size=7.5, color=DGRAY)
    T(sl, exp,     xs2[1], y, ws2[1], 0.28, size=7.5, color=DGRAY, align=PP_ALIGN.CENTER)
    T(sl, iv,      xs2[2], y, ws2[2], 0.28, size=7.5, color=DGRAY, align=PP_ALIGN.CENTER)
    T(sl, verdict, xs2[3], y, ws2[3], 0.28, size=7.5, bold=True, color=vc)
    y += 0.34

# Recommendation box below chart
R(sl, 0.45, 5.35, 7.9, 1.7, LTEAL, line_color=TEAL)
R(sl, 0.45, 5.35, 0.06, 1.7, TEAL)
T(sl, 'Proposed Weight Rebalancing', 0.65, 5.43, 7.5, 0.3,
  size=9.5, bold=True, color=NAVY)
T(sl, '↑ Increase C3 → 30%  |  ↑ F1, F4, F5 → 15% each  |  ↓ C1/C2 → 10% each  |  ✕ Remove F3 (IV < 0.05)',
  0.65, 5.78, 7.55, 0.38, size=9, color=NAVY)
T(sl, 'Rationale: IV directly measures the log-odds separating power of each feature — a more objective basis than expert judgement alone.',
  0.65, 6.18, 7.55, 0.68, size=8.5, color=DGRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Q2: Scorecard Enhancements
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'III. Scorecard Enhancement Suggestions',
               'Six targeted improvements to model robustness, calibration and predictive power')
footer(sl, 8)

enhancements = [
    ('1. Laplace Smoothing on WoE', TEAL,
     ['Sparse bins produce WoE = ±20, distorting all scores',
      'Apply ε = 0.5: pct_bad = (bad + ε) / (total_bad + ε × n)',
      'Prevents extreme score swings for near-zero event bins']),
    ('2. Bin Merging & Monotonicity', TEAL,
     ['"Restructured in past" has 0 defaults — merge with adjacent',
      'Enforce monotone WoE: merge if direction breaks risk order',
      'Rule: merge bins with n < 5 or zero events without risk logic']),
    ('3. Remove F3 Sale Growth', RED,
     ['IV = 0.04 → below weak threshold of 0.10',
      'Chi-square p = 0.89 — not statistically significant',
      'Adds noise without predictive gain — omit from scorecard']),
    ('4. Recalibrate PD Probabilities', AMBER,
     ['class_weight="balanced" inflates predicted PD 4×',
      'Apply Platt scaling (CalibratedClassifierCV, cv=5)',
      'Or recalculate base_odds from actual portfolio default rate']),
    ('5. Macroeconomic Overlay', NAVY,
     ['Add GDP growth, interest rate, oil price as macro variables',
      'Test 1–4 quarter lags via Granger causality',
      'Two-stage model: micro scorecard × macro scalar adjustment']),
    ('6. Rating Scale Calibration', NAVY,
     ['Use log-linear PD spacing (PD doubles each rating grade)',
      'Each band needs ≥ 20 defaults for stable PD estimates',
      'Validate: PD must strictly increase from best to worst band']),
]
xs3 = [0.45, 4.55, 8.65]
ys3 = [1.28, 4.08]
idx = 0
for r in range(2):
    for c in range(3):
        title, accent, bullets = enhancements[idx]
        card_top(sl, xs3[c], ys3[r], 3.9, 2.55, title, bullets, accent=accent)
        idx += 1

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Q3: EDA Methodology
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'II. EDA Methodology Framework',
               'Generalised five-step approach for credit risk datasets with borrower, financial and macro data')
footer(sl, 9)

steps = [
    ('Step 1', 'Data Quality', NAVY,
     ['Missing value audit & sentinel detection',
      'Skewness check → IQR-based capping',
      'Data type and duplicate validation']),
    ('Step 2', 'Univariate Analysis', TEAL,
     ['Categorical: frequency, default rate per bin',
      'Numerical: histogram, boxplot, skewness',
      'Flag >50% single-bin concentration']),
    ('Step 3', 'Bivariate Analysis', NAVY,
     ['Numerical: KS test, Pearson correlation',
      'Categorical: Chi-square, Cramér\'s V',
      'IV/WoE: quantify per-bin predictive power']),
    ('Step 4', 'Multivariate', TEAL,
     ['Pearson + Cramér\'s V combined matrix',
      'VIF for numerical multicollinearity',
      'PCA if > 15 features (dimensionality)']),
    ('Step 5', 'Macro Integration', NAVY,
     ['GDP growth, rates, oil vs default rate',
      'Granger causality with 1–4Q lags',
      'Two-stage: micro model + macro scalar']),
]
x = 0.45
for step, title, color, bullets in steps:
    R(sl, x, 1.28, 2.44, 5.75, WHITE, line_color=LGRAY)
    R(sl, x, 1.28, 2.44, 0.85, color)
    T(sl, step,  x+0.1, 1.32, 2.24, 0.3,  size=8, color=GOLD, align=PP_ALIGN.CENTER)
    T(sl, title, x+0.1, 1.62, 2.24, 0.48, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    divider(sl, 2.14, l=x+0.1, w=2.24, color=LGRAY, thick=0.02)
    y = 2.22
    for b in bullets:
        T(sl, f'• {b}', x+0.18, y, 2.1, 0.42, size=8.5, color=DGRAY)
        y += 0.55
    if steps.index((step, title, color, bullets)) < 4:
        R(sl, x+2.44, 2.75, 0.11, 0.25, GOLD)  # arrow connector
    x += 2.57

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Q4: PD Model Development
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'IV. PD Model Development — Logistic Regression',
               'Theory, WoE encoding rationale, assumptions and PDO score calibration')
footer(sl, 10)

# Left — theory
R(sl, 0.45, 1.28, 6.1, 5.75, WHITE, line_color=LGRAY)
R(sl, 0.45, 1.28, 6.1, 0.42, NAVY)
T(sl, 'Model Theory & Assumptions', 0.65, 1.33, 5.8, 0.32, size=10, bold=True, color=WHITE)
theory = [
    ('Log-Odds Formulation',
     'log(p / 1−p) = β₀ + β₁WoE₁ + … + βₙWoEₙ\nWoE encoding linearises the log-odds relationship for LR'),
    ('MLE Estimation',
     'Parameters found by maximising log-likelihood (no closed form)\nSolver: iterative Newton-Raphson / L-BFGS'),
    ('WoE Encoding Effect',
     'Replaces bins with ln(%Bad / %Good) — continuous, monotone\nHandles non-linearity; nullifies need for dummy variables'),
    ('Class Imbalance',
     'class_weight="balanced" re-weights loss function per class\nNote: inflates predicted PD — recalibration step mandatory'),
    ('Key Assumptions',
     'Linear log-odds · No perfect multicollinearity\nObservations i.i.d. · Large-sample MLE validity'),
]
y = 1.83
for title, detail in theory:
    T(sl, f'▸ {title}', 0.65, y,      5.7, 0.25, size=8.5, bold=True, color=TEAL)
    T(sl, detail,        0.65, y+0.25, 5.7, 0.46, size=8,   color=DGRAY)
    y += 0.80

# Right — implementation steps
R(sl, 6.75, 1.28, 6.15, 5.75, WHITE, line_color=LGRAY)
R(sl, 6.75, 1.28, 6.15, 0.42, TEAL)
T(sl, 'Implementation Pipeline', 6.95, 1.33, 5.85, 0.32, size=10, bold=True, color=WHITE)
steps6 = [
    ('Data Preparation',
     'Bin features; fill NaN → "Not available"; cap sentinel values'),
    ('WoE Encoding',
     'Compute WoE per bin with Laplace smoothing ε = 0.5'),
    ('Train / Test Split',
     '80/20 stratified; EPV check: need ≥ 10 defaults per feature'),
    ('Model Training',
     'LogisticRegression(class_weight="balanced"); no scaling needed'),
    ('PDO Calibration',
     'Factor = PDO/ln2 = 28.85 · Offset = 600 − 28.85×ln(50) = 487.1'),
    ('Score Mapping',
     'Score = Offset + Factor×β₀ + Σ[−Factor × βᵢ × WoEᵢ]'),
]
y = 1.83
for i, (title, detail) in enumerate(steps6):
    R(sl, 6.78, y-0.02, 0.42, 0.42, TEAL if i % 2 == 0 else NAVY)
    T(sl, str(i+1), 6.80, y,   0.38, 0.32, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, title,    7.30, y,   5.45, 0.25, size=8.5, bold=True, color=TEAL)
    T(sl, detail,   7.30, y+0.25, 5.45, 0.42, size=8, color=DGRAY)
    y += 0.80

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Model Results
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'IV. Model Results',
               'LR coefficients, normalised importance vs expert weights and performance metrics')
footer(sl, 11)

# KPI strip
kpi_data = [
    ('0.90', 'AUC-ROC',    'Excellent', TEAL,  TEAL),
    ('0.80', 'Gini',       'Strong',    TEAL,  TEAL),
    ('65%',  'Accuracy',   'Threshold-\ndependent', NAVY, NAVY),
    ('80%',  'Recall (D)', 'Default\ndetection',    TEAL, TEAL),
    ('28.85\n/ 487.1', 'Factor / Offset', 'PDO=20, Base=600', NAVY, NAVY),
]
x = 0.45
for val, lbl, sub, vc, tc in kpi_data:
    kpi(sl, x, 1.28, 2.44, 1.2, val, lbl, sub, vc, tc)
    x += 2.55

# Coef table
R(sl, 0.45, 2.65, 5.55, 4.55, WHITE, line_color=LGRAY)
R(sl, 0.45, 2.65, 5.55, 0.42, NAVY)
T(sl, 'Coefficients & Feature Importance', 0.65, 2.70, 5.25, 0.32, size=9.5, bold=True, color=WHITE)
coef_rows = [
    ('C3 CB Delinquency',  '0.897','26.6%','25%',  GREEN, LGREEN),
    ('F4 D/E Ratio',       '0.745','22.2%','10%',  RED,   LRED),
    ('F5 DSCR',            '0.545','16.2%','10%',  RED,   LRED),
    ('C2 Debt Rescheduling','0.469','14.0%','12.5%',AMBER, LGOLD),
    ('F1 Net Fixed Asset', '0.282',' 8.4%','10%',  AMBER, LGOLD),
    ('C1 Debt Restructuring','0.238',' 7.1%','12.5%',AMBER,LGOLD),
    ('F2 Net Profit Margin','0.182',' 5.5%','10%',  AMBER, LGOLD),
]
xc = [0.50, 2.35, 3.38, 4.3]
wc = [1.80,  0.98, 0.87, 0.85]
hd = ['Feature', 'Coef', 'LR Wt%', 'Expert']
y = 3.18
for i, h in enumerate(hd):
    T(sl, h, xc[i], y, wc[i], 0.27, size=8, bold=True, color=NAVY)
divider(sl, 3.48, l=0.5, w=5.4, color=TEAL)
y = 3.55
for feat, coef, lrw, exp, vc, bgc in coef_rows:
    R(sl, 0.47, y-0.02, 5.49, 0.31, bgc)
    T(sl, feat, xc[0], y, wc[0], 0.28, size=7.5, color=DGRAY)
    T(sl, coef, xc[1], y, wc[1], 0.28, size=7.5, color=DGRAY, align=PP_ALIGN.CENTER)
    T(sl, lrw,  xc[2], y, wc[2], 0.28, size=7.5, bold=True, color=vc, align=PP_ALIGN.CENTER)
    T(sl, exp,  xc[3], y, wc[3], 0.28, size=7.5, color=DGRAY, align=PP_ALIGN.CENTER)
    y += 0.32

# charts
PIC(sl, 'roc_curve_chart.png',      6.15, 2.65, 3.55, 2.25)
PIC(sl, 'coef_vs_expert_chart.png', 9.85, 2.65, 3.4,  2.25)
PIC(sl, 'score_band_chart.png',     6.15, 5.0,  7.1,  2.2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Q5: Model Evaluation
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'V. Model Validation & Evaluation',
               'Three pillars: goodness-of-fit, discriminatory power and forecasting accuracy')
footer(sl, 12)

pillars = [
    (NAVY, 'Pillar 1\nGoodness-of-Fit', [
        ('Hosmer-Lemeshow Test',
         'Stat = 297.94  ·  p ≈ 0.00 → Poor calibration fit\n'
         'Root cause: extreme WoE (±20) from zero-event bins\n'
         'Fix: Laplace smoothing + bin merging before retest'),
        ('Binomial Test',
         'Observed: 50 defaults (9.0%)  ·  Predicted avg PD: 34.97%\n'
         'Model over-predicts PD by 4× — due to class_weight="balanced"\n'
         'Fix: Platt scaling or recalibrate base_odds to actual rate'),
    ]),
    (TEAL, 'Pillar 2\nDiscriminatory Power', [
        ('AUC-ROC = 0.90  ·  Gini = 0.80',
         'Model correctly ranks 90% of default/non-default pairs\n'
         'Excellent discrimination even with miscalibrated probabilities\n'
         'AUC is robust to class imbalance — reliable metric here'),
        ('KS Statistic',
         'Score means: 395.9 (default) vs 636.6 (non-default)\n'
         'Strong 240-point separation between risk groups\n'
         'Top score band (300–498): 44.8% default rate'),
    ]),
    (GOLD, 'Pillar 3\nForecasting Accuracy', [
        ('PSI — Population Stability Index',
         'Monitor score distribution shift on new data over time\n'
         'PSI < 0.10: stable · 0.10–0.25: investigate · > 0.25: rebuild\n'
         'Run monthly after deployment'),
        ('Back-Testing & Validation',
         'Compare predicted PD vs realised default rates per score band\n'
         'Validate annually; flag bands where |PD_pred − PD_actual| > 2%\n'
         'Maintain champion/challenger framework for ongoing monitoring'),
    ]),
]
x = 0.45
for color, title, items in pillars:
    R(sl, x, 1.28, 4.1, 5.75, WHITE, line_color=LGRAY)
    R(sl, x, 1.28, 4.1, 0.75, color)
    T(sl, title, x+0.15, 1.3, 3.82, 0.65, size=11, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    y = 2.15
    for metric, detail in items:
        R(sl, x+0.15, y, 3.82, 0.3, LGRAY)
        T(sl, metric, x+0.22, y+0.04, 3.68, 0.24, size=8.5, bold=True, color=NAVY)
        T(sl, detail, x+0.22, y+0.36, 3.68, 0.8,  size=8,   color=DGRAY)
        y += 1.25
    x += 4.3

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Score Distribution
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl); header(sl, 'V. Score Distribution & Band Analysis',
               'PDO-calibrated scores scaled to 300–850 with band-level default rate validation')
footer(sl, 13)

PIC(sl, 'score_band_chart.png', 0.45, 1.28, 8.7, 4.45)

R(sl, 9.35, 1.28, 3.58, 4.45, WHITE, line_color=LGRAY)
R(sl, 9.35, 1.28, 3.58, 0.42, NAVY)
T(sl, 'PDO Parameters', 9.55, 1.33, 3.3, 0.32, size=10, bold=True, color=WHITE)
pdo_params = [
    ('PDO',        '20 pts  (odds double per 20 pts)'),
    ('Base Score', '600  at  50:1 odds'),
    ('Factor',     'PDO / ln(2)  =  28.85'),
    ('Offset',     '600 − 28.85 × ln(50)  =  487.1'),
    ('Intercept',  'Factor × β₀  =  10.55'),
    ('Scale',      '300–850  (min-max normalised)'),
]
y = 1.83
for k, v in pdo_params:
    bg3 = LGRAY if pdo_params.index((k,v)) % 2 == 0 else WHITE
    R(sl, 9.37, y-0.02, 3.52, 0.34, bg3)
    T(sl, k, 9.50, y+0.01, 1.1,  0.28, size=8, bold=True, color=TEAL)
    T(sl, v, 10.65, y+0.01, 2.22, 0.28, size=8, color=DGRAY)
    y += 0.36

divider(sl, 3.85, l=9.38, w=3.5, color=TEAL)
T(sl, 'Band Findings', 9.55, 3.93, 3.3, 0.28, size=9, bold=True, color=NAVY)
band_obs = [
    ('300–498', '44.8% default', 'HIGH RISK',   RED),
    ('498–513', '5–13% default', 'MEDIUM RISK', AMBER),
    ('513–850', '0–3.5% default','LOW RISK',    GREEN),
]
y = 4.25
for band, dr, label, lc in band_obs:
    R(sl, 9.37, y, 3.52, 0.42, LRED if lc == RED else (LGOLD if lc == AMBER else LGREEN))
    T(sl, band,  9.50, y+0.07, 0.85, 0.28, size=8, bold=True, color=NAVY)
    T(sl, dr,    10.38, y+0.07, 1.3,  0.28, size=8, color=DGRAY)
    T(sl, label, 11.72, y+0.07, 1.05, 0.28, size=7.5, bold=True, color=lc)
    y += 0.47

divider(sl, 5.85)
R(sl, 0.45, 5.9, 12.43, 1.0, LRED, line_color=RED)
R(sl, 0.45, 5.9, 0.06, 1.0, RED)
T(sl, '⚠  Calibration Issue',
  0.62, 5.96, 3.0, 0.28, size=9, bold=True, color=RED)
T(sl, 'Predicted avg PD = 34.97% vs Actual 9.0% — class_weight="balanced" inflates probabilities 4×. '
       'Apply Platt scaling (CalibratedClassifierCV) before using PD values for provisioning, pricing or regulatory reporting.',
  0.62, 6.27, 12.1, 0.55, size=8.5, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion & Recommendations
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
R(sl, 0, 0, 13.33, 7.5, NAVY2)
R(sl, 0, 0, 0.42,  7.5, GOLD)
R(sl, 0.42, 0,    12.91, 0.07, GOLD)
R(sl, 0.42, 7.43, 12.91, 0.07, TEAL)

T(sl, 'Conclusion & Recommendations', 1.0, 0.2, 12.0, 0.72,
  size=24, bold=True, color=WHITE)
divider(sl, 1.05, l=1.0, w=11.85, color=GOLD, thick=0.05)

recs_final = [
    (GOLD, 'Rebalance Feature Weights',
     'Raise C3 → 30%, D/E Ratio & DSCR → 15% each.\nRemove F3 Sale Growth (IV = 0.04).'),
    (GOLD, 'Fix WoE Sparse Bins',
     'Apply Laplace smoothing ε = 0.5 to all bins.\nMerge bins with n < 5 or zero-default cells.'),
    (GOLD, 'Recalibrate PD',
     'class_weight distorts PD 4×.\nApply Platt scaling or use actual portfolio base_odds.'),
    (TEAL, 'Validate Robustly',
     'Re-run HL test after bin fixes.\nImplement PSI monitoring for ongoing score drift.'),
    (TEAL, 'Add Macro Variables',
     'Integrate GDP growth, interest rate, oil price.\nTwo-stage: micro scorecard + macro scalar.'),
    (TEAL, 'Rating Scale Design',
     'Log-linear PD spacing (PD doubles per grade).\n≥ 20 defaults per band for stable estimates.'),
]
xs4 = [0.55, 4.65, 8.75]
ys4 = [1.22, 3.92]
idx = 0
for r in range(2):
    for c in range(3):
        color, title, detail = recs_final[idx]
        R(sl, xs4[c], ys4[r], 3.9, 2.48, RGBColor(0x0E, 0x30, 0x5F))
        R(sl, xs4[c], ys4[r], 3.9, 0.07, color)
        R(sl, xs4[c], ys4[r], 0.07, 2.48, color)
        T(sl, title,  xs4[c]+0.2, ys4[r]+0.12, 3.55, 0.35,
          size=10, bold=True, color=WHITE)
        T(sl, detail, xs4[c]+0.2, ys4[r]+0.52, 3.55, 1.8,
          size=9, color=MGRAY)
        idx += 1

R(sl, 0.55, 6.55, 12.33, 0.65, RGBColor(0x0E, 0x30, 0x5F))
T(sl, 'Model achieves AUC = 0.90 — strong discrimination — but requires calibration fixes before production deployment.',
  0.75, 6.63, 12.0, 0.48, size=10, color=GOLD, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = '/Users/leekim/prj/deloitte/Credit_Scorecard_Presentation_v4.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
